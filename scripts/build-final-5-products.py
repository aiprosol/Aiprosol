#!/usr/bin/env python3
"""
Aiprosol — Final 5 product deliverables:
  • Enterprise Readiness Kit ($397) → XLSX scorecard + DOCX RFP + PPTX 90-day plan
  • Zapier + Make Bundle ($197)     → 50 recipes × 2 platforms (markdown)
  • ROI Calculator ($47)            → XLSX with working formulas
  • ROI Pitch Deck ($17)            → PPTX template
  • Audit Checklist ($37)           → Interactive HTML
"""

from pathlib import Path
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.shapes import MSO_SHAPE
import json

CATALOG = Path("/Users/user/Airprosol/products catalogue/01-ready-to-sell")

# Aiprosol brand colours
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
LIGHT_VIOLET = RGBColor(0xC0, 0x84, 0xFC)
DARK_TEXT = RGBColor(0x1F, 0x1B, 0x2E)
MUTED = RGBColor(0x6B, 0x75, 0x85)
BG = RGBColor(0xFA, 0xF8, 0xFF)

VIOLET_HEX = "8B5CF6"
LIGHT_VIOLET_HEX = "C084FC"
DARK_TEXT_HEX = "1F1B2E"
MUTED_HEX = "6B7585"
BG_HEX = "FAF8FF"
EDGE_HEX = "EDE9FE"

PVIOLET = PRGBColor(0x8B, 0x5C, 0xF6)
PLVIOLET = PRGBColor(0xC0, 0x84, 0xFC)
PDARK = PRGBColor(0x1F, 0x1B, 0x2E)
PMUTED = PRGBColor(0x6B, 0x75, 0x85)


# ════════════════════════════════════════════════════════════════════
# 1. ENTERPRISE READINESS KIT — XLSX + DOCX + PPTX
# ════════════════════════════════════════════════════════════════════

ER_DIR = CATALOG / "assessments/397-enterprise-ai-readiness-assessment-kit/delivery"
ER_DIR.mkdir(parents=True, exist_ok=True)

# ─── 1a) Readiness Scorecard XLSX (120 questions, 12 dimensions) ───
DIMENSIONS = [
    ("Data Hygiene", [
        "Is customer data deduplicated across systems?",
        "Are key entities (customer, account, deal) defined consistently?",
        "Is there a single source of truth for each critical entity?",
        "Is data freshness monitored (last-modified visibility)?",
        "Are PII fields catalogued and access-logged?",
        "Are data quality issues tracked in a ticket queue?",
        "Is there a defined data steward role?",
        "Are schema changes communicated before deployment?",
        "Is historical data preserved (vs overwritten) for analysis?",
        "Are data warehouses refreshed on a predictable cadence?",
    ]),
    ("Process Documentation", [
        "Are top 20 business processes documented in writing?",
        "Are SOPs reviewed at least annually?",
        "Can a new hire follow an SOP without senior help?",
        "Are exceptions (vs happy path) documented?",
        "Is there a process owner for each documented workflow?",
        "Are process changes version-controlled?",
        "Are bottlenecks identified and tracked?",
        "Are cycle times measured for key processes?",
        "Are upstream/downstream dependencies mapped?",
        "Are documented processes actually followed (audited)?",
    ]),
    ("Change Management", [
        "Is there a defined change-approval process?",
        "Are stakeholders identified before changes?",
        "Is there a communication plan for major changes?",
        "Are training resources created for new tools?",
        "Are rollback procedures documented?",
        "Are change-impact assessments performed?",
        "Is feedback collected post-implementation?",
        "Are champions identified within each business unit?",
        "Are resistance patterns recognised and addressed?",
        "Is change saturation monitored (avoiding fatigue)?",
    ]),
    ("Governance", [
        "Is there an AI governance policy?",
        "Are data privacy regulations (GDPR/CCPA/etc) documented?",
        "Are AI use cases reviewed before deployment?",
        "Is there a model risk management framework?",
        "Are bias audits performed on AI outputs?",
        "Is there a compliance officer or equivalent?",
        "Are vendor risk assessments performed?",
        "Are audit trails maintained for AI decisions?",
        "Is there an AI ethics board or review process?",
        "Are policies updated as regulations evolve?",
    ]),
    ("Integration Debt", [
        "Are integrations documented in a system catalog?",
        "Are API contracts version-controlled?",
        "Are deprecated integrations actively retired?",
        "Are point-to-point integrations limited (vs hub model)?",
        "Are integration failures monitored and alerted?",
        "Are retries and idempotency standard?",
        "Are credentials rotated on a schedule?",
        "Are webhooks signed and verified?",
        "Are rate limits respected with backoff?",
        "Are integration costs tracked per workflow?",
    ]),
    ("Talent & Skills", [
        "Are AI literacy programs available to all staff?",
        "Are technical leads assessed on AI capability?",
        "Is there a clear career path for data/AI roles?",
        "Are external AI experts engaged for specialist work?",
        "Is upskilling budget allocated annually?",
        "Are AI champions identified outside IT?",
        "Are job descriptions updated to reflect AI competencies?",
        "Are interview processes assessing AI judgment?",
        "Are AI ethics included in employee training?",
        "Is there a community of practice for AI users?",
    ]),
    ("Vendor Strategy", [
        "Is there a vendor selection framework?",
        "Are concentration risks managed (multi-vendor)?",
        "Are contracts reviewed for AI-specific clauses?",
        "Are exit costs evaluated before signing?",
        "Is vendor performance tracked against SLAs?",
        "Are vendor security postures audited?",
        "Are data residency requirements specified?",
        "Are vendor-locked features identified?",
        "Are open-source alternatives evaluated?",
        "Is total cost of ownership (TCO) computed?",
    ]),
    ("Security", [
        "Are AI workloads in a separate security zone?",
        "Is data masked before being sent to external AI?",
        "Are prompts logged and reviewed for sensitive data?",
        "Are outputs validated for data leakage?",
        "Are access controls based on principle of least privilege?",
        "Are AI model weights protected from extraction?",
        "Are jailbreak attempts monitored?",
        "Are incident response plans tested?",
        "Are penetration tests performed on AI endpoints?",
        "Are security patches applied within SLA?",
    ]),
    ("Finance Operations", [
        "Are AI/automation costs tracked at workflow level?",
        "Is ROI measured per automation?",
        "Are budgets reviewed quarterly?",
        "Are cost spikes alerted automatically?",
        "Are cost-saving wins documented and shared?",
        "Are licenses reconciled against actual usage?",
        "Are spot-savings opportunities reviewed (e.g., spot instances)?",
        "Are TCO models maintained for major systems?",
        "Are unit economics calculated for each product line?",
        "Are CFO dashboards updated weekly?",
    ]),
    ("Customer Operations", [
        "Are customer journey maps documented?",
        "Are NPS surveys deployed at key touchpoints?",
        "Are support tickets tagged for product feedback?",
        "Are customer health scores computed?",
        "Are churn signals tracked weekly?",
        "Are expansion signals tracked weekly?",
        "Are customer success playbooks documented?",
        "Are voice-of-customer themes surfaced monthly?",
        "Is the customer data platform deduplicated?",
        "Are loyalty programs measured for retention impact?",
    ]),
    ("Product Engineering", [
        "Are feature deployments automated (CI/CD)?",
        "Are feature flags used for safe rollouts?",
        "Are A/B tests performed before major releases?",
        "Are technical debt tickets tracked?",
        "Is code review universal?",
        "Are observability dashboards owner-tagged?",
        "Are post-deploy verifications automated?",
        "Are uptime SLAs published and tracked?",
        "Are AI features marked clearly in product?",
        "Are user research insights tied to roadmap?",
    ]),
    ("Strategic Alignment", [
        "Is the CEO articulating an AI strategy publicly?",
        "Are AI initiatives tied to top-3 company OKRs?",
        "Is there a 12-month AI roadmap?",
        "Are AI investments reviewed at board level?",
        "Is competitive intelligence informed by AI capability?",
        "Are M&A targets evaluated for AI capability?",
        "Is brand positioning updated for AI era?",
        "Are partnerships sought to fill AI gaps?",
        "Are customer expectations of AI tracked?",
        "Is the company telling its AI story externally?",
    ]),
]


def build_scorecard():
    wb = Workbook()

    # Cover sheet
    cover = wb.active
    cover.title = "Cover"
    cover["A1"] = "Aiprosol Enterprise AI Readiness Assessment"
    cover["A1"].font = Font(name="Calibri", size=24, bold=True, color=VIOLET_HEX)
    cover["A3"] = "120 questions across 12 dimensions"
    cover["A3"].font = Font(name="Calibri", size=14, color=MUTED_HEX)
    cover["A5"] = "How to use:"
    cover["A5"].font = Font(name="Calibri", size=12, bold=True)
    instructions = [
        "1. On the Scorecard tab, answer each question 0-5 (5 = mature, 0 = absent).",
        "2. The Summary tab auto-aggregates your score per dimension (0-100).",
        "3. The Heatmap tab highlights your weakest dimensions in red.",
        "4. The Action Map tab maps each gap to a 12-month roadmap recommendation.",
        "5. Re-run quarterly to track maturity progress.",
        "",
        "Scoring guide:",
        "   5 — Fully implemented, regularly audited",
        "   4 — Implemented, occasionally audited",
        "   3 — Partial / inconsistent",
        "   2 — Documented intent, not yet implemented",
        "   1 — Awareness, no plan",
        "   0 — Not on the radar",
    ]
    for i, line in enumerate(instructions, start=6):
        cover[f"A{i}"] = line
        cover[f"A{i}"].font = Font(name="Calibri", size=11)
    cover.column_dimensions["A"].width = 100

    # Scorecard sheet
    sc = wb.create_sheet("Scorecard")
    sc.append(["Dimension", "Question", "Score (0-5)", "Notes"])
    for col_idx, header in enumerate(["Dimension", "Question", "Score (0-5)", "Notes"], start=1):
        cell = sc.cell(row=1, column=col_idx)
        cell.font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=VIOLET_HEX)
        cell.alignment = Alignment(horizontal="center", vertical="center")

    row = 2
    for dim_name, questions in DIMENSIONS:
        for q in questions:
            sc.cell(row=row, column=1, value=dim_name)
            sc.cell(row=row, column=2, value=q)
            sc.cell(row=row, column=3, value=0)  # default
            sc.cell(row=row, column=4, value="")
            row += 1

    # column widths
    sc.column_dimensions["A"].width = 28
    sc.column_dimensions["B"].width = 80
    sc.column_dimensions["C"].width = 12
    sc.column_dimensions["D"].width = 50

    # Conditional formatting on Score column
    rule = ColorScaleRule(start_type="num", start_value=0, start_color="EF4444",
                          mid_type="num", mid_value=3, mid_color="F59E0B",
                          end_type="num", end_value=5, end_color="22C55E")
    sc.conditional_formatting.add(f"C2:C{row-1}", rule)

    # Summary sheet — aggregate per dimension
    summary = wb.create_sheet("Summary")
    summary.append(["Dimension", "Score (0-100)", "Status"])
    for col_idx in range(1, 4):
        cell = summary.cell(row=1, column=col_idx)
        cell.font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=VIOLET_HEX)

    start_row = 2
    cur_q = 2
    for dim_name, questions in DIMENSIONS:
        end_q = cur_q + len(questions) - 1
        summary.cell(row=start_row, column=1, value=dim_name)
        summary.cell(row=start_row, column=2, value=f"=ROUND(AVERAGE(Scorecard!C{cur_q}:C{end_q})*20, 0)")
        summary.cell(row=start_row, column=3, value=f'=IF(B{start_row}>=80,"Mature",IF(B{start_row}>=60,"Capable",IF(B{start_row}>=40,"Emerging","Critical")))')
        cur_q = end_q + 1
        start_row += 1

    summary.cell(row=start_row, column=1, value="OVERALL").font = Font(bold=True)
    summary.cell(row=start_row, column=2,
                 value=f"=ROUND(AVERAGE(B2:B{start_row-1}), 0)").font = Font(bold=True, color=VIOLET_HEX)
    summary.column_dimensions["A"].width = 28
    summary.column_dimensions["B"].width = 18
    summary.column_dimensions["C"].width = 18

    rule2 = ColorScaleRule(start_type="num", start_value=0, start_color="EF4444",
                           mid_type="num", mid_value=50, mid_color="F59E0B",
                           end_type="num", end_value=100, end_color="22C55E")
    summary.conditional_formatting.add(f"B2:B{start_row}", rule2)

    # Add a bar chart
    chart = BarChart()
    chart.type = "bar"
    chart.title = "Score by Dimension"
    data_ref = Reference(summary, min_col=2, min_row=1, max_col=2, max_row=start_row-1)
    cats_ref = Reference(summary, min_col=1, min_row=2, max_row=start_row-1)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)
    chart.height = 12
    chart.width = 18
    summary.add_chart(chart, "E2")

    # Action Map sheet
    am = wb.create_sheet("Action Map")
    am.append(["Dimension", "Score band", "Recommended 12-month action"])
    for col_idx in range(1, 4):
        cell = am.cell(row=1, column=col_idx)
        cell.font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=VIOLET_HEX)

    ACTIONS = {
        "Critical (0-39)": "Halt new AI initiatives. Run a 90-day foundation sprint: hire a dimension owner, document baseline, ship one quick win.",
        "Emerging (40-59)": "Pick 1 quick win + 1 structural fix in Q1. Begin quarterly review cadence.",
        "Capable (60-79)": "Push for differentiation. Identify 2 advanced use cases that compound.",
        "Mature (80-100)": "Maintain + share patterns externally. Look for opportunities to commercialise expertise."
    }
    row = 2
    for dim_name, _ in DIMENSIONS:
        for band, action in ACTIONS.items():
            am.cell(row=row, column=1, value=dim_name)
            am.cell(row=row, column=2, value=band)
            am.cell(row=row, column=3, value=action)
            row += 1
    am.column_dimensions["A"].width = 28
    am.column_dimensions["B"].width = 22
    am.column_dimensions["C"].width = 90

    out = ER_DIR / "01-readiness-scorecard.xlsx"
    wb.save(out)
    return out


# ─── 1b) RFP template DOCX ───
def build_rfp():
    doc = Document()

    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    title = doc.add_heading("RFP Template — AI Automation Vendor Selection", level=1)
    for run in title.runs:
        run.font.color.rgb = VIOLET
        run.font.size = Pt(20)
    doc.add_paragraph("Aiprosol Enterprise AI Readiness Assessment Kit · 2026")

    doc.add_heading("1. About this RFP", level=2)
    doc.add_paragraph(
        "This RFP template guides procurement teams through selecting AI automation vendors. "
        "It emphasises operational readiness, data governance, and exit costs — not feature lists."
    )

    doc.add_heading("2. Buyer profile", level=2)
    table = doc.add_table(rows=6, cols=2)
    table.style = "Light Grid Accent 4"
    fields = [
        ("Organisation", "{{COMPANY NAME}}"),
        ("Industry / sector", "{{INDUSTRY}}"),
        ("Annual revenue band", "{{REVENUE BAND}}"),
        ("Employees", "{{HEADCOUNT}}"),
        ("Current AI maturity (overall score)", "{{SCORE}}/100"),
        ("Strategic AI sponsor", "{{NAME, ROLE}}"),
    ]
    for i, (k, v) in enumerate(fields):
        table.rows[i].cells[0].text = k
        table.rows[i].cells[1].text = v

    doc.add_heading("3. The problem we are solving", level=2)
    doc.add_paragraph(
        "{{Describe the specific operational problem, including: current cost (hours/$), "
        "scale (transactions/month), pain owner, why it matters now, why it has not been solved internally.}}"
    )

    doc.add_heading("4. What we need", level=2)
    doc.add_paragraph("Must-haves:")
    must_haves = [
        "Native integrations with our current stack (list below)",
        "On-prem or sovereign-cloud deployment option",
        "Audit logging of all AI decisions",
        "PII redaction before any data leaves our perimeter",
        "Defined exit / data-export procedure",
        "Reference customers in our industry willing to be contacted",
    ]
    for m in must_haves:
        doc.add_paragraph(m, style="List Bullet")
    doc.add_paragraph("Nice-to-haves:")
    nice = [
        "Pre-built workflows for our use case",
        "Bring-your-own-model support",
        "SAML / SCIM",
        "Granular feature toggles",
        "European data residency",
    ]
    for n in nice:
        doc.add_paragraph(n, style="List Bullet")

    doc.add_heading("5. Current stack to integrate with", level=2)
    table2 = doc.add_table(rows=6, cols=2)
    table2.style = "Light Grid Accent 4"
    table2.rows[0].cells[0].text = "Layer"
    table2.rows[0].cells[1].text = "Tool"
    table2.rows[0].cells[0].paragraphs[0].runs[0].font.bold = True
    table2.rows[0].cells[1].paragraphs[0].runs[0].font.bold = True
    for i, (k, v) in enumerate([
        ("CRM", "{{HubSpot/Salesforce/...}}"),
        ("ERP", "{{NetSuite/SAP/...}}"),
        ("Ticketing", "{{Linear/Jira/...}}"),
        ("Data warehouse", "{{Snowflake/BigQuery/...}}"),
        ("Comms", "{{Slack/Teams/...}}"),
    ], start=1):
        table2.rows[i].cells[0].text = k
        table2.rows[i].cells[1].text = v

    doc.add_heading("6. Evaluation criteria (40-criteria grid)", level=2)
    criteria = [
        ("Capability", "Does the product solve the stated problem end-to-end?"),
        ("Capability", "Does it handle our specific edge cases (list 3)?"),
        ("Capability", "Does it offer the 'must-have' integrations?"),
        ("Capability", "Are AI outputs auditable post-hoc?"),
        ("Capability", "Can humans intervene at any step?"),
        ("Security", "SOC 2 Type II report available?"),
        ("Security", "PII handling documented?"),
        ("Security", "Encryption at rest + in transit?"),
        ("Security", "Pen-test report (last 12 months)?"),
        ("Security", "Sub-processor list disclosed?"),
        ("Reliability", "Uptime SLA with penalties?"),
        ("Reliability", "Incident response SLA defined?"),
        ("Reliability", "Status page with historical data?"),
        ("Reliability", "Disaster recovery tested?"),
        ("Cost", "Total monthly cost at expected scale?"),
        ("Cost", "Cost predictability (no usage-based surprises)?"),
        ("Cost", "Cost of running for 3 years vs alternative?"),
        ("Cost", "Hidden costs (egress, support tiers)?"),
        ("Implementation", "Time to first value (weeks)?"),
        ("Implementation", "Professional services available?"),
        ("Implementation", "Customer success model defined?"),
        ("Implementation", "Onboarding playbook documented?"),
        ("Exit", "Data export format documented?"),
        ("Exit", "Time to fully export your data?"),
        ("Exit", "Termination notice period acceptable?"),
        ("Exit", "Migration tools for switching?"),
        ("Vendor health", "Funding runway > 24 months?"),
        ("Vendor health", "Customer concentration risk acceptable?"),
        ("Vendor health", "Leadership team stable?"),
        ("Vendor health", "Public roadmap aligned with our needs?"),
        ("Support", "Response time SLA?"),
        ("Support", "Dedicated CSM at our spend level?"),
        ("Support", "24×7 vs business hours?"),
        ("Support", "Knowledge base quality?"),
        ("References", "3 customers we can talk to in our industry"),
        ("References", "1 reference at our scale"),
        ("References", "1 reference that has been a customer >24 months"),
        ("References", "1 reference that has churned (and reason)"),
        ("Future-fit", "Roadmap visibility for next 12 months?"),
        ("Future-fit", "Ability to bring our own model?"),
    ]
    table3 = doc.add_table(rows=1, cols=3)
    table3.style = "Light Grid Accent 4"
    table3.rows[0].cells[0].text = "Category"
    table3.rows[0].cells[1].text = "Criterion"
    table3.rows[0].cells[2].text = "Score (1-5)"
    for cat, crit in criteria:
        row = table3.add_row()
        row.cells[0].text = cat
        row.cells[1].text = crit
        row.cells[2].text = "___"

    doc.add_heading("7. Contract red-flag checklist", level=2)
    redflags = [
        "Auto-renewal without notice obligation",
        "Unilateral price-increase clause without cap",
        "Data ownership ambiguity",
        "Mandatory arbitration in unfavourable jurisdiction",
        "Unilateral change-of-terms",
        "IP rights to your data / configurations",
        "Liability cap below 1× annual fees",
        "Lack of force-majeure carve-outs for security incidents",
        "Vague SLA enforcement mechanism",
        "Termination penalty greater than 1 quarter of fees",
    ]
    for r in redflags:
        doc.add_paragraph(r, style="List Bullet")

    doc.add_heading("8. Response format we expect", level=2)
    doc.add_paragraph(
        "Vendors should provide: (1) a 1-page executive summary, "
        "(2) the criteria grid completed with scores + evidence, "
        "(3) all-in pricing for 3 years at our expected scale, "
        "(4) references contact details, and (5) any deviations from this RFP listed explicitly."
    )

    doc.add_heading("9. Timeline", level=2)
    timeline = [
        ("Issue RFP", "Week 0"),
        ("Vendor questions", "Week 1"),
        ("Q&A response", "Week 2"),
        ("Responses due", "Week 3"),
        ("Shortlist demos", "Week 4"),
        ("Reference calls", "Week 5"),
        ("Final selection", "Week 6"),
        ("Contract execution", "Week 7-8"),
    ]
    table4 = doc.add_table(rows=1, cols=2)
    table4.style = "Light Grid Accent 4"
    for k, v in timeline:
        r = table4.add_row()
        r.cells[0].text = k
        r.cells[1].text = v

    doc.add_paragraph()
    doc.add_paragraph(
        "© 2026 Aiprosol Ltd · Enterprise AI Readiness Assessment Kit · Licensed to purchaser for internal use."
    ).runs[0].font.color.rgb = MUTED

    out = ER_DIR / "02-rfp-vendor-selection.docx"
    doc.save(out)
    return out


# ─── 1c) 90-day plan PPTX ───
def build_90day_plan():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    def add_slide(title, subtitle=None, bullets=None, layout=5):
        slide = prs.slides.add_slide(prs.slide_layouts[layout])

        # Aiprosol-style background accent
        tb = slide.shapes.add_textbox(PInches(0.5), PInches(0.3), PInches(12.3), PInches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.runs[0].font.bold = True
        p.runs[0].font.size = PPt(34)
        p.runs[0].font.color.rgb = PVIOLET
        p.runs[0].font.name = "Calibri"

        if subtitle:
            sb = slide.shapes.add_textbox(PInches(0.5), PInches(1.1), PInches(12.3), PInches(0.5))
            ps = sb.text_frame.paragraphs[0]
            ps.text = subtitle
            ps.runs[0].font.size = PPt(16)
            ps.runs[0].font.color.rgb = PMUTED
            ps.runs[0].font.name = "Calibri"

        if bullets:
            body = slide.shapes.add_textbox(PInches(0.5), PInches(1.8), PInches(12.3), PInches(5.2))
            tf = body.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {b}"
                p.font.size = PPt(18)
                p.font.color.rgb = PDARK
                p.font.name = "Calibri"
                p.space_after = PPt(8)

        # Footer
        f = slide.shapes.add_textbox(PInches(0.5), PInches(7.0), PInches(12.3), PInches(0.4))
        fp = f.text_frame.paragraphs[0]
        fp.text = "Aiprosol Enterprise AI Readiness · 90-Day Implementation Plan"
        fp.runs[0].font.size = PPt(10)
        fp.runs[0].font.color.rgb = PMUTED
        return slide

    # Cover
    cover = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    title_box = cover.shapes.add_textbox(PInches(1), PInches(2.5), PInches(11), PInches(1.5))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "90-Day AI Readiness Implementation Plan"
    tp.runs[0].font.size = PPt(40)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = PVIOLET
    sub = cover.shapes.add_textbox(PInches(1), PInches(4), PInches(11), PInches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "From assessment to first measurable outcome"
    sp.runs[0].font.size = PPt(22)
    sp.runs[0].font.color.rgb = PMUTED
    foot = cover.shapes.add_textbox(PInches(1), PInches(6.5), PInches(11), PInches(0.5))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "© 2026 Aiprosol Ltd · Enterprise AI Readiness Assessment Kit"
    fp.runs[0].font.size = PPt(11)
    fp.runs[0].font.color.rgb = PMUTED

    add_slide("Three phases", "Foundation → Pilot → Scale", [
        "Phase 0 (Days 1-30) — Foundation: alignment, governance, baseline measurement",
        "Phase 1 (Days 31-60) — Pilot: 1 high-leverage use case, end-to-end in production",
        "Phase 2 (Days 61-90) — Scale: 3 more use cases, repeatable playbook"
    ])

    add_slide("Phase 0 · Foundation (Days 1-30)",
              "Get alignment + measure baseline",
              [
                  "Day 1-7: Stakeholder alignment (CEO + CIO + line-of-business head)",
                  "Day 8-14: Run the 120-question Readiness Scorecard with affected teams",
                  "Day 15-21: Identify top 3 candidate use cases via Opportunity Map",
                  "Day 22-28: Select 1 use case for pilot (highest ROI × lowest risk)",
                  "Day 29-30: Sign-off on success criteria + check-in cadence",
              ])
    add_slide("Phase 0 · Roles", None, [
        "Executive sponsor — accountable for outcome",
        "Pilot lead — runs daily standup, manages timeline",
        "Domain expert — has the operational context",
        "Technical lead — implements the workflow",
        "Change champion — manages stakeholder comms",
    ])
    add_slide("Phase 0 · Deliverables", None, [
        "Completed Readiness Scorecard (all 120 questions)",
        "Opportunity Map (top 5 use cases ranked by value × feasibility)",
        "Signed-off pilot charter (use case, success criteria, timeline, budget)",
        "Vendor RFP (if external tooling needed)",
        "Communications plan",
    ])

    add_slide("Phase 1 · Pilot (Days 31-60)",
              "Ship one use case end-to-end",
              [
                  "Day 31-37: Set up tooling + integrations",
                  "Day 38-49: Build the workflow with humans-in-the-loop (Phase 1 of the 3-phase progression)",
                  "Day 50-56: Move to Phase 2 (auto with override) once quality is proven",
                  "Day 57-60: Measure baseline vs after; document learnings",
              ])
    add_slide("Phase 1 · Success criteria", None, [
        "Quantitative outcome achieved (e.g., hours saved / week)",
        "Quality maintained or improved (vs manual baseline)",
        "Zero security or compliance incidents",
        "Adoption by intended users >80%",
        "Total time-to-value documented for replication",
    ])
    add_slide("Phase 1 · Risks to manage", None, [
        "Scope creep — adding 'just one more' use case before the first works",
        "Quality regression — moving to auto before quality is proven",
        "Adoption resistance — solve with champion + visible wins",
        "Cost surprises — set up cost monitoring before launch",
        "Vendor lock-in — keep an exit plan in mind from Day 1",
    ])

    add_slide("Phase 2 · Scale (Days 61-90)",
              "Repeat with 3 more use cases + build muscle memory",
              [
                  "Day 61-67: Document pilot's playbook (the repeatable steps)",
                  "Day 68-74: Begin use cases #2-#3 in parallel (different teams, similar pattern)",
                  "Day 75-81: Begin use case #4 (start to test cross-system patterns)",
                  "Day 82-90: Quarterly business review — what to scale next, what to retire",
              ])
    add_slide("Phase 2 · Compounding", None, [
        "Cross-workflow data: signals from one workflow feed another",
        "Shared infrastructure: common LLM endpoints, prompt library, vector store",
        "Shared learnings: weekly retro across pilot teams",
        "Shared metrics: same dashboards across use cases",
    ])

    add_slide("90-day outcomes you should expect", None, [
        "4 working AI-augmented workflows in production",
        "Documented playbook for the next 10",
        "Measured baseline + improvement for each",
        "Team aligned + confident on next phase",
        "Defensible business case for year-2 investment",
    ])
    add_slide("Year-2 horizon (preview)", None, [
        "Move from 4 → 30 workflows with the same team size",
        "Begin internal AI centre of excellence",
        "Take learnings to executive education + board",
        "Look for category-defining moves (not just optimisation)",
    ])

    add_slide("Common 90-day pitfalls", None, [
        "Trying to fix governance + pilot a use case + train teams all in Day 1",
        "Picking too big a pilot ('let's automate all of sales')",
        "Selecting tools before defining the problem",
        "Underinvesting in change management",
        "Declaring victory before measuring outcomes",
    ])

    add_slide("Aiprosol contact", "When you need expert help",
              ["hello@aiprosol.com", "aiprosol.com", "Srijan Paudel, Founder"])

    out = ER_DIR / "03-90-day-implementation-plan.pptx"
    prs.save(out)
    return out


er_files = [build_scorecard(), build_rfp(), build_90day_plan()]
print(f"✓ Enterprise Readiness Kit: {len(er_files)} files")
for f in er_files:
    print(f"  {f.stat().st_size/1024:.1f} KB  {f.name}")


# ════════════════════════════════════════════════════════════════════
# 2. ZAPIER + MAKE BUNDLE — 50 build-step recipes per platform
# ════════════════════════════════════════════════════════════════════

ZM_DIR = CATALOG / "playbooks/197-zapier-make-power-user-bundle/delivery"
ZM_DIR.mkdir(parents=True, exist_ok=True)

ZAPIER_RECIPES = [
    ("Stripe payment → CRM contact + Slack",
     ["1. Trigger: Stripe — New Payment Intent Succeeded",
      "2. Action: Filter — Only continue if amount > 100",
      "3. Action: Formatter (Text) — Extract customer first name from email",
      "4. Action: HubSpot — Create or Update Contact",
      "5. Action: HubSpot — Create Deal stage='closed-won' amount=Stripe amount",
      "6. Action: Slack — Send Channel Message #sales-wins"]),
    ("Calendly booking → CRM + prep brief Slack DM",
     ["1. Trigger: Calendly — New Event Created",
      "2. Action: HubSpot — Find Contact (by email)",
      "3. Filter: If contact doesn't exist, create",
      "4. Action: Webhook GET — fetch company info from clearbit",
      "5. Action: OpenAI — Create Chat Completion for 3-bullet prep brief",
      "6. Action: Slack — DM rep with brief + calendar link"]),
    ("Form submission → AI score → Tier routing",
     ["1. Trigger: Webhooks — Catch Hook (form POST)",
      "2. Action: OpenAI — Chat Completion to score 0-100",
      "3. Action: Formatter (Numbers) — Parse score from text",
      "4. Action: Paths — Branch by tier: hot (>80) / warm (60-80) / cold",
      "5. Hot path: Slack #sales + Cal.com priority booking link",
      "6. Warm path: Add to Mailchimp 'nurture' segment",
      "7. Cold path: Add to Mailchimp 'newsletter only' segment"]),
    ("Gmail label 'INVOICE' → OCR → Sheet log",
     ["1. Trigger: Gmail — New Email Matching Search 'label:invoice'",
      "2. Action: Filter — Has Attachment = true",
      "3. Action: OpenAI Vision — Extract {vendor, amount, due_date}",
      "4. Action: Formatter — Parse JSON response",
      "5. Action: Google Sheets — Create Spreadsheet Row (Accounting Log)",
      "6. Action: Google Drive — Upload File (archive PDF)"]),
    ("HubSpot deal stage change → Email draft",
     ["1. Trigger: HubSpot — Deal Stage Updated",
      "2. Action: HubSpot — Get associated primary contact",
      "3. Action: OpenAI — Chat Completion for next-step email based on new stage",
      "4. Action: Gmail — Create Draft (NOT send) for review"]),
]

# Use first 5 as fully-detailed examples; add 45 more as 1-line recipes
MORE_ZAPIER_RECIPES_TITLES = [
    "Typeform → Airtable + Slack",
    "Shopify order → Klaviyo segment",
    "Stripe failed payment → Slack + dunning email",
    "Intercom chat → AI summary → Notion",
    "Google Forms intake → Zendesk ticket",
    "ClickUp overdue task → Slack DM assignee",
    "Mailchimp campaign sent → Sheet log",
    "Notion page update → Slack thread",
    "Calendly cancellation → CRM property update",
    "Salesforce opportunity → Slack #revenue",
    "Trello card moved → Slack channel ping",
    "Webflow CMS publish → Sheet log + tweet",
    "GitHub PR opened → Slack #eng",
    "Stripe subscription cancelled → CS alert",
    "QuickBooks invoice paid → CRM update",
    "Asana task completed → Manager email",
    "Twilio SMS → AI reply → CRM note",
    "Eventbrite registration → AI welcome email",
    "Discord new member → AI welcome DM",
    "RSS feed update → Tweet draft",
    "Notion DB add → Calendar event create",
    "Reddit mention → Slack ping",
    "Twitter mention → CRM activity log",
    "Stripe refund → CS investigation ticket",
    "Calendly meeting end → Linear task",
    "HubSpot form fill → AI follow-up email draft",
    "Stripe subscription renewed → thank-you email",
    "Slack reaction added → save to Notion",
    "Google Drive new file → Sheet log",
    "Outlook new email → AI classify + label",
    "Square sale → revenue Sheet log",
    "Pipedrive deal won → onboarding kickoff workflow",
    "GitHub issue → Linear sync",
    "Loom recording shared → AI transcript + summary",
    "Notion comment → email author",
    "Salesforce lead update → Slack notification",
    "Asana new task → AI breakdown into subtasks",
    "Stripe customer churn signal → CS alert",
    "Calendly multi-attendee booking → Notion meeting prep",
    "Twilio call missed → CRM activity log",
    "Trello label change → Slack DM",
    "Shopify abandoned cart → Klaviyo recovery flow",
    "Mailchimp unsubscribe → CRM property update",
    "GitHub release → social media post draft",
    "Stripe disputed charge → finance alert",
]

MAKE_RECIPES = [
    ("HubSpot lead → enrichment → routing",
     ["1. Trigger: HubSpot · Watch contacts (created or updated)",
      "2. Module: HTTP · Make a Request to Clearbit /v2/combined/find",
      "3. Module: OpenAI · Create Completion to compute lead score",
      "4. Module: Router with 3 paths: hot, warm, cold",
      "5. Hot: Slack · Send a Message to #sales + HubSpot · Update Contact",
      "6. Warm: Mailchimp · Add Member to List + HubSpot · Update Contact",
      "7. Cold: HubSpot · Update Contact only",
      "8. Error handler: dead-letter to Google Sheets"]),
    ("Stripe payment → multi-step CRM + accounting",
     ["1. Trigger: Stripe · Watch Events filtered to payment_intent.succeeded",
      "2. Module: HubSpot · Search Contacts by email",
      "3. Router: contact exists / new contact",
      "4. Module: HubSpot · Create Deal (existing) OR Create Contact then Deal (new)",
      "5. Module: Xero · Create Invoice (status=PAID)",
      "6. Module: Slack · Send a Message #revenue",
      "7. Iterator: For each line item → Sheet log"]),
    ("Daily 8am cron → pipeline aggregation → Slack",
     ["1. Trigger: Schedule · 0 8 * * 1-5",
      "2. Module: HubSpot · Search deals (last 7 days)",
      "3. Module: Aggregator · Group by deal stage",
      "4. Module: Tools · Set Multiple Variables (totals per stage)",
      "5. Module: OpenAI · Chat Completion for 3-bullet narrative",
      "6. Module: Slack · Send a Message #sales-leadership",
      "7. Module: Google Sheets · Add a Row to weekly archive"]),
    ("Meeting transcript → action items → Linear",
     ["1. Trigger: Webhook · Catch Hook from Fireflies",
      "2. Module: OpenAI · Chat Completion to extract action items as JSON",
      "3. Module: Tools · Parse JSON",
      "4. Module: Iterator · Loop over action items",
      "5. Module: Linear · Create Issue per item",
      "6. Module: Slack · Send a Message with summary count"]),
    ("Drive new file → AI category → folder routing",
     ["1. Trigger: Google Drive · Watch Files in Folder (Inbox)",
      "2. Module: Google Drive · Download a File",
      "3. Module: OpenAI Vision · Classify document",
      "4. Module: Router by category: invoice/contract/receipt/other",
      "5. Each path: Google Drive · Move a File to destination folder",
      "6. Module: Slack · Send a Message #automations-log"]),
]

MORE_MAKE_RECIPES_TITLES = [
    "Typeform → Airtable with iterators",
    "Shopify order → multi-vendor split routing",
    "Salesforce opportunity stage → orchestrator with delays",
    "Stripe webhook → idempotent CRM upsert",
    "Notion DB → cross-DB sync with state store",
    "GitHub PR → Slack + Linear with router",
    "Mailchimp campaign metrics → BigQuery daily ETL",
    "Calendly bulk bookings → AI prep briefs in parallel",
    "Customer support tickets → AI classify → router",
    "Stripe MRR daily aggregation",
    "Multi-source data merge with data stores",
    "Long-running customer onboarding orchestrator",
    "Idempotent webhook processor with dedupe",
    "Rate-limited API consumer with backoff",
    "Conditional approval gate (Slack buttons)",
    "Cross-org Notion → Google Workspace sync",
    "Slack thread reply → external API webhook",
    "Twilio voicemail → AI transcript → ticket",
    "Cron-driven QBR document assembly",
    "Vendor renewal radar (60-day forward)",
    "Async standup aggregator (DMs in parallel)",
    "Document approval workflow with audit log",
    "Multi-step e-commerce post-purchase",
    "AI content moderation pipeline",
    "Self-healing flow with retry policies",
    "Cost monitoring + spike alerts",
    "Workflow health check meta-automation",
    "Customer success signals aggregator",
    "Multi-channel publishing fan-out",
    "Survey response → AI sentiment → segments",
    "Inventory sync across Shopify + WMS",
    "Lead-to-meeting scheduling chain",
    "Renewal forecast workflow",
    "ChatGPT response auto-improver",
    "GitHub release → multi-channel announcement",
    "Daily error log aggregator",
    "AI-assisted refund decision workflow",
    "Customer churn prediction pipeline",
    "Pricing change notification cascade",
    "Compliance audit trail recorder",
    "Multi-language customer support router",
    "Dynamic content personalization workflow",
    "Vendor invoice multi-approval chain",
    "Recruiting pipeline orchestrator",
    "AI-driven A/B test analysis",
    "Sales territory rebalancer",
]


def write_recipes(filename, platform_name, full_recipes, more_titles):
    lines = [f"# {platform_name} Power User — 50 Production Recipes", "",
             f"Version 1.0 · 50 recipes · © 2026 Aiprosol Ltd", "",
             "## Format", "",
             "First 5 recipes are FULL recipes with every step.",
             "Recipes 6-50 are 1-line patterns — apply the templates from recipes 1-5.",
             "", "## Full recipes (1-5)", ""]
    for i, (title, steps) in enumerate(full_recipes, 1):
        lines.append(f"### {i}. {title}")
        for s in steps:
            lines.append(f"   {s}")
        lines.append("")
    lines.append("## Pattern list (6-50)")
    lines.append("")
    for i, t in enumerate(more_titles, 6):
        lines.append(f"{i}. {t}")
    lines.append("")
    lines.append("## How to apply")
    lines.append("Each title maps to one of the 7 patterns from the Aiprosol Workflow Playbook:")
    lines.append("- Linear pipeline · Branching by classifier · Fan-out · Scheduled aggregation")
    lines.append("- Polling-with-state · Approval gate · Long-running orchestrator")
    lines.append("")
    lines.append("Use the first 5 full recipes as templates — substitute trigger, intermediate steps, and destination.")
    out = ZM_DIR / filename
    out.write_text("\n".join(lines))
    return out


zapier_md = write_recipes("zapier-50-recipes.md", "Zapier", ZAPIER_RECIPES, MORE_ZAPIER_RECIPES_TITLES)
make_md = write_recipes("make-50-recipes.md", "Make.com", MAKE_RECIPES, MORE_MAKE_RECIPES_TITLES)

# Decision matrix
decision_md = ZM_DIR / "zapier-vs-make-decision-matrix.md"
decision_md.write_text("""# Zapier vs Make — Decision Matrix

Aiprosol Power User Bundle · © 2026

## When to pick Zapier
- 6000+ integrations needed — Zapier has the broadest catalog
- Speed-to-first-automation matters — Zapier is faster to build simple zaps
- Polish + reliability over flexibility
- Non-technical team building automations
- Per-task pricing acceptable for low/medium volume

## When to pick Make
- Complex workflows with iterators, aggregators, data stores
- Higher volume — Make's pricing scales better at 10k+ ops/month
- Visual scenario design helps non-technical reviewers
- Need data transformations beyond Zapier Formatter
- Open to learning curve in exchange for power

## When to skip both and use n8n
- Need self-host for security / compliance
- Want unlimited tasks at flat pricing
- Open-source preferences
- Engineering team available to manage

## 22-row comparison

| Dimension | Zapier | Make |
|---|---|---|
| Integrations count | ~6000 | ~2000 |
| Pricing starts | $20/mo | $9/mo |
| Free tier | 100 tasks/mo | 1000 ops/mo |
| Editor | Linear, simpler | Visual graph, more powerful |
| Iterators | Limited | First-class |
| Aggregators | No | Yes |
| Data stores | No | Yes |
| Error handling | Basic | Advanced |
| Webhooks | Yes | Yes (advanced) |
| HTTP requests | Yes | Yes |
| Code steps | Yes (limited) | Yes |
| AI integrations | Excellent | Very good |
| Team workflows | Excellent | Good |
| Version control | Yes (paid) | Yes |
| Audit log | Yes (paid) | Yes |
| Multi-step paths | Yes (paid) | Yes (free) |
| Sub-zaps / sub-scenarios | Yes (paid) | Yes |
| Scheduled triggers | Yes | Yes |
| Polling triggers | Yes | Yes |
| Real-time webhook | Yes | Yes |
| Self-host option | No | No (n8n if you need it) |
| Support tier | Strong | Solid |
""")

print(f"✓ Zapier + Make Bundle: 3 markdown files")


# ════════════════════════════════════════════════════════════════════
# 3. ROI CALCULATOR — XLSX with formulas
# ════════════════════════════════════════════════════════════════════

RC_DIR = CATALOG / "templates/047-ai-automation-roi-calculator/delivery"
RC_DIR.mkdir(parents=True, exist_ok=True)


def build_roi_calculator():
    wb = Workbook()

    # Inputs sheet
    inp = wb.active
    inp.title = "Inputs"
    inp["A1"] = "AI Automation ROI Calculator"
    inp["A1"].font = Font(name="Calibri", size=22, bold=True, color=VIOLET_HEX)
    inp["A2"] = "Aiprosol · 2026 · Fill in B-column cells in yellow"
    inp["A2"].font = Font(name="Calibri", size=11, color=MUTED_HEX)

    inputs = [
        ("VOLUME", None, None),
        ("Transactions per month (current)", 1000, "Count of items currently processed"),
        ("Transactions per month (projected with automation)", 1000, "Will it stay the same or grow?"),
        ("", None, None),
        ("TIME", None, None),
        ("Average time per transaction — manual (minutes)", 5, "How long currently"),
        ("Average time per transaction — automated (minutes)", 0.5, "How long after"),
        ("Working hours per month", 160, "Standard ~160 = 40 × 4 weeks"),
        ("", None, None),
        ("PEOPLE COSTS", None, None),
        ("Fully-loaded labour cost per hour (USD)", 50, "Salary × 1.4 typically, ÷ working hours"),
        ("Number of people doing this work currently", 2, "Headcount touching this process"),
        ("", None, None),
        ("IMPLEMENTATION COSTS", None, None),
        ("One-time implementation cost (USD)", 5000, "Build + onboarding"),
        ("Monthly tool cost (USD)", 50, "n8n + OpenAI + ancillary tools"),
        ("Monthly maintenance hours", 4, "Hours/month to monitor + fix"),
        ("", None, None),
        ("CONFIDENCE", None, None),
        ("Quality vs manual (1=worse, 1.0=same, >1=better)", 1.0, "Adjust if AI improves quality"),
        ("Adoption ramp curve — month 1", 0.3, "Fraction of full benefit in month 1"),
        ("Adoption ramp curve — month 3", 0.8, "Fraction of full benefit by month 3"),
        ("Adoption ramp curve — month 6+", 1.0, "Fraction of full benefit at maturity"),
        ("Risk-adjusted confidence (0-1)", 0.8, "Apply this haircut to projected savings"),
    ]

    row = 4
    for label, default, note in inputs:
        if default is None and label and not note:
            # Section header
            inp.cell(row=row, column=1, value=label).font = Font(name="Calibri", size=12, bold=True, color=VIOLET_HEX)
            inp.cell(row=row, column=1).fill = PatternFill("solid", fgColor=EDGE_HEX)
        elif label:
            inp.cell(row=row, column=1, value=label).font = Font(name="Calibri", size=11)
            cell = inp.cell(row=row, column=2, value=default)
            cell.font = Font(name="Calibri", size=11, bold=True)
            cell.fill = PatternFill("solid", fgColor="FEF3C7")  # input cell — yellow
            if note:
                inp.cell(row=row, column=3, value=note).font = Font(name="Calibri", size=10, color=MUTED_HEX, italic=True)
        row += 1

    inp.column_dimensions["A"].width = 50
    inp.column_dimensions["B"].width = 15
    inp.column_dimensions["C"].width = 60

    # Calc sheet
    calc = wb.create_sheet("Calculations")
    calc["A1"] = "Auto-computed — do not edit"
    calc["A1"].font = Font(name="Calibri", size=14, bold=True, color=VIOLET_HEX)

    rows = [
        ("Hours saved per month at maturity",
         "=Inputs!B7*(Inputs!B10-Inputs!B11)/60"),
        ("Monthly $ saved at maturity (gross)",
         "=B3*Inputs!B14"),
        ("Monthly maintenance cost ($)",
         "=Inputs!B22*Inputs!B14"),
        ("Net monthly savings at maturity ($)",
         "=B4-B5-Inputs!B21"),
        ("Risk-adjusted net monthly savings ($)",
         "=B6*Inputs!B30"),
        ("", ""),
        ("Implementation cost ($)",
         "=Inputs!B20"),
        ("Payback period at maturity (months)",
         "=IF(B7>0, B9/B7, \"N/A\")"),
        ("", ""),
        ("12-month gross savings ($) — accounting for ramp",
         "=B7*(Inputs!B27*3 + Inputs!B28*3 + Inputs!B29*6) - B9"),
        ("12-month savings at NPV (8% discount)",
         "=B12 / 1.08"),
        ("", ""),
        ("ROI multiple over 12 months",
         "=IF(B9>0, B12/B9, \"N/A\")"),
        ("Hours/year reclaimed",
         "=B3*12"),
    ]
    for i, (label, formula) in enumerate(rows, start=3):
        calc.cell(row=i, column=1, value=label).font = Font(name="Calibri", size=11)
        if formula:
            cell = calc.cell(row=i, column=2)
            if formula.startswith("="):
                cell.value = formula
            cell.font = Font(name="Calibri", size=11, bold=True, color=VIOLET_HEX)
    calc.column_dimensions["A"].width = 50
    calc.column_dimensions["B"].width = 20

    # Monthly projection sheet
    proj = wb.create_sheet("12-Month Projection")
    proj["A1"] = "12-Month Cash Flow Projection"
    proj["A1"].font = Font(name="Calibri", size=14, bold=True, color=VIOLET_HEX)
    proj.append([])
    proj.append(["Month", "Adoption %", "Gross $ saved", "Maintenance $", "Net $ saved", "Cumulative net"])
    for col_idx in range(1, 7):
        cell = proj.cell(row=3, column=col_idx)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=VIOLET_HEX)

    for m in range(1, 13):
        # Adoption curve: simple step function
        if m <= 3:
            adoption = "=Inputs!B27"
        elif m <= 6:
            adoption = "=Inputs!B28"
        else:
            adoption = "=Inputs!B29"
        proj.cell(row=3 + m, column=1, value=m)
        proj.cell(row=3 + m, column=2, value=adoption)
        proj.cell(row=3 + m, column=3, value=f"=B{3+m}*Calculations!B4")
        proj.cell(row=3 + m, column=4, value=f"=Calculations!B5")
        proj.cell(row=3 + m, column=5, value=f"=C{3+m}-D{3+m}-IF({m}=1, Inputs!B20, 0)")
        if m == 1:
            proj.cell(row=3 + m, column=6, value=f"=E{3+m}")
        else:
            proj.cell(row=3 + m, column=6, value=f"=F{3+m-1}+E{3+m}")
    for col in "BCDEF":
        proj.column_dimensions[col].width = 16

    # Chart
    chart = LineChart()
    chart.title = "Cumulative Net Savings ($)"
    chart.style = 12
    chart.x_axis.title = "Month"
    chart.y_axis.title = "USD"
    data = Reference(proj, min_col=6, min_row=3, max_row=15)
    cats = Reference(proj, min_col=1, min_row=4, max_row=15)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 10
    chart.width = 16
    proj.add_chart(chart, "H3")

    # Output sheet
    out_ws = wb.create_sheet("Output (1-page)")
    out_ws["A1"] = "ROI Summary"
    out_ws["A1"].font = Font(size=22, bold=True, color=VIOLET_HEX)
    out_ws["A3"] = "Monthly net savings"
    out_ws["B3"] = "=Calculations!B6"
    out_ws["A4"] = "Payback period"
    out_ws["B4"] = "=Calculations!B10 & \" months\""
    out_ws["A5"] = "12-month gross savings"
    out_ws["B5"] = "=Calculations!B12"
    out_ws["A6"] = "ROI multiple (12mo)"
    out_ws["B6"] = "=Calculations!B14 & \"x\""
    out_ws["A7"] = "Hours reclaimed / year"
    out_ws["B7"] = "=Calculations!B15"
    for r in range(3, 8):
        out_ws.cell(row=r, column=1).font = Font(size=12, bold=True)
        out_ws.cell(row=r, column=2).font = Font(size=14, color=VIOLET_HEX)
    out_ws.column_dimensions["A"].width = 30
    out_ws.column_dimensions["B"].width = 30

    out = RC_DIR / "ai-automation-roi-calculator.xlsx"
    wb.save(out)
    return out


roi_xlsx = build_roi_calculator()
print(f"✓ ROI Calculator: {roi_xlsx.stat().st_size/1024:.1f} KB  {roi_xlsx.name}")


# ════════════════════════════════════════════════════════════════════
# 4. ROI PITCH DECK — PPTX template (25 slides)
# ════════════════════════════════════════════════════════════════════

PD_DIR = CATALOG / "templates/017-automation-roi-pitch-deck-template/delivery"
PD_DIR.mkdir(parents=True, exist_ok=True)


def build_pitch_deck():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    def add(title, subtitle=None, bullets=None, notes=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(PInches(0.6), PInches(0.4), PInches(12.1), PInches(0.9))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = title
        tp.runs[0].font.size = PPt(32)
        tp.runs[0].font.bold = True
        tp.runs[0].font.color.rgb = PVIOLET
        tp.runs[0].font.name = "Calibri"

        if subtitle:
            sb = slide.shapes.add_textbox(PInches(0.6), PInches(1.3), PInches(12.1), PInches(0.5))
            ps = sb.text_frame.paragraphs[0]
            ps.text = subtitle
            ps.runs[0].font.size = PPt(15)
            ps.runs[0].font.color.rgb = PMUTED

        if bullets:
            body = slide.shapes.add_textbox(PInches(0.6), PInches(2.0), PInches(12.1), PInches(4.8))
            tf = body.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {b}"
                p.font.size = PPt(18)
                p.font.color.rgb = PDARK
                p.font.name = "Calibri"
                p.space_after = PPt(6)

        # Notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        # Footer
        f = slide.shapes.add_textbox(PInches(0.6), PInches(6.9), PInches(12.1), PInches(0.4))
        fp = f.text_frame.paragraphs[0]
        fp.text = "{{Your Company}} · Automation ROI Pitch · 2026"
        fp.runs[0].font.size = PPt(10)
        fp.runs[0].font.color.rgb = PMUTED
        return slide

    # Slide 1: Cover (custom)
    cov = prs.slides.add_slide(prs.slide_layouts[6])
    t = cov.shapes.add_textbox(PInches(1), PInches(2.8), PInches(11.3), PInches(1.5))
    tp = t.text_frame.paragraphs[0]
    tp.text = "Automating {{Process Name}}"
    tp.runs[0].font.size = PPt(44)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = PVIOLET
    s = cov.shapes.add_textbox(PInches(1), PInches(4.3), PInches(11.3), PInches(0.8))
    sp = s.text_frame.paragraphs[0]
    sp.text = "${{Annual Savings}} reclaimed · {{Payback}}-month payback · {{Hours}} hrs/yr freed"
    sp.runs[0].font.size = PPt(22)
    sp.runs[0].font.color.rgb = PMUTED

    f2 = cov.shapes.add_textbox(PInches(1), PInches(6.5), PInches(11.3), PInches(0.5))
    fp2 = f2.text_frame.paragraphs[0]
    fp2.text = "Prepared for {{Audience}} · {{Date}}"
    fp2.runs[0].font.size = PPt(13)
    fp2.runs[0].font.color.rgb = PMUTED

    # Slides 2-25
    add("Agenda", None, [
        "1. The problem we're solving",
        "2. Current cost of inaction",
        "3. Proposed approach",
        "4. Projected savings",
        "5. Payback + risk analysis",
        "6. Implementation plan",
        "7. The ask",
    ], notes="Establishes time expectations and signals which audiences should focus on which slides.")

    add("The problem", "{{1-sentence problem statement}}", [
        "Who feels it: {{Persona affected}}",
        "How often: {{Frequency}}",
        "Hidden cost: {{Cost they don't notice}}",
        "Why now: {{Trigger event making this urgent}}",
    ], notes="Frame the problem in the audience's language. Make it specific. Q: 'Hasn't this always been a problem?' A: 'Yes — what changed is {{trigger event}}.'")

    add("Current state — by the numbers", "Where the money is going", [
        "Manual time per transaction: {{minutes}} min",
        "Transactions per month: {{count}}",
        "Monthly hours consumed: {{hours}}",
        "Annual cost (loaded): ${{annual_cost}}",
        "Quality issues observed: {{count}}",
    ], notes="Use the ROI Calculator output to populate. Anchor on concrete numbers.")

    add("Who feels it the most", None, [
        "{{Team A}}: {{specific pain}}",
        "{{Team B}}: {{specific pain}}",
        "{{Customer impact}}: {{description}}",
        "{{Internal impact on hiring/retention}}: {{description}}",
    ], notes="Mention specific people by role. Q: 'Why haven't they fixed it themselves?' A: 'They've tried — see appendix slide 22.'")

    add("Proposed approach", None, [
        "Phase 1 (weeks 1-2): {{Step}}",
        "Phase 2 (weeks 3-6): {{Step}}",
        "Phase 3 (weeks 7-12): {{Step}}",
        "Human-in-the-loop until quality is proven (Phase 1 = AI-as-suggestion)",
        "Move to auto-action only with {{measurable confidence threshold}}",
    ], notes="Phased approach helps de-risk. Q: 'Why can't we go straight to auto?' A: 'Because we'd lose the ability to catch errors before they reach customers.'")

    add("How it works (1-min explanation)", None, [
        "Input: {{What triggers the automation}}",
        "Step 1: {{First processing step}}",
        "Step 2: {{Second step}}",
        "Step 3: {{Decision or output}}",
        "Output: {{What the user sees}}",
        "Fallback: {{What happens if it fails}}",
    ], notes="Use plain language. The CFO doesn't need to understand the LLM API.")

    add("Projected savings", "Year 1 economics", [
        "Hours saved / year: {{hours}}",
        "Direct cost savings: ${{savings}}",
        "Quality improvement: {{quantified outcome}}",
        "Capacity unlocked: {{what people can do instead}}",
        "Customer outcome: {{what improves for them}}",
    ], notes="Connect the savings to a customer outcome where possible. CFOs find that compelling.")

    add("Payback analysis", None, [
        "One-time implementation cost: ${{implementation_cost}}",
        "Monthly maintenance: ${{maintenance}}",
        "Monthly net savings: ${{net_savings}}",
        "Payback period: {{months}} months",
        "12-month ROI multiple: {{multiple}}x",
    ], notes="Pull these numbers from the ROI Calculator. Anchor on payback period — it's the number CFOs internalise.")

    add("Sensitivity analysis", "What if our assumptions are wrong?", [
        "If volume is 50% lower: payback = {{X}} months",
        "If quality regression hurts customers: cost = ${{Y}}",
        "If labour costs change: net savings = ${{Z}}",
        "If implementation overruns 50%: payback = {{X}} months",
    ], notes="Pre-empt the 'but what if' questions. Q: 'These assumptions seem optimistic.' A: 'Sensitivity table — even under conservative assumptions, payback < {{X}} months.'")

    add("Case study — {{Industry}}", "How {{Similar Company}} did this", [
        "Problem: {{their problem}}",
        "Approach: {{what they did}}",
        "Result: {{their numbers}}",
        "Quote: '{{quote from their team}}'",
    ], notes="Substitute with a real reference. Aiprosol provides case studies for Legal/RealEstate/Manufacturing.")

    add("Case study — {{Industry 2}}", None, [
        "Problem: {{their problem}}",
        "Result: {{their numbers}}",
    ], notes="Second case study slot.")

    add("Case study — {{Industry 3}}", None, [
        "Problem: {{their problem}}",
        "Result: {{their numbers}}",
    ], notes="Third case study slot.")

    add("Implementation plan — 90 days", None, [
        "Days 1-7: Alignment + tooling setup",
        "Days 8-21: Build with humans-in-the-loop",
        "Days 22-49: Pilot in one team",
        "Days 50-90: Scale to {{N}} teams",
        "Day 90: Quarterly business review",
    ], notes="Show realistic pacing. Q: 'Can you do it in 30 days?' A: 'We can build the workflow in 14 days. The remaining 76 are about adoption and confidence.'")

    add("Team + roles", None, [
        "Executive sponsor: {{Name}}",
        "Project lead: {{Name}}",
        "Domain expert: {{Name}}",
        "Technical lead: {{Name}} (or vendor)",
        "Change champion: {{Name}}",
    ], notes="Named people > 'a team will'. Forces accountability.")

    add("Risk register", "The 5 things that could go wrong", [
        "Quality regression — mitigation: HITL phase 1-2",
        "Adoption resistance — mitigation: champion + visible wins",
        "Cost overrun — mitigation: weekly cost monitoring",
        "Integration breakage — mitigation: monitoring + dead letter queue",
        "Vendor lock-in — mitigation: exit plan defined Day 1",
    ], notes="Don't pretend risks don't exist. Stating them + mitigations builds trust.")

    add("Vendor + tooling decisions", None, [
        "AI model: {{OpenAI / Claude / open-source}} (rationale: {{...}})",
        "Workflow engine: {{n8n / Zapier / Make}}",
        "CRM integration: {{HubSpot / Salesforce / ...}}",
        "Monitoring: {{tool}}",
        "Build vs buy: {{decision + rationale}}",
    ], notes="Pre-empt the 'why this tool' question. Q: 'Why not Microsoft / our existing vendor?' A: '{{specific rationale}}'")

    add("Security + compliance", None, [
        "Data handling: {{encryption, residency, retention}}",
        "Access control: {{who can see what}}",
        "Audit trail: {{retention period}}",
        "Compliance: {{GDPR / SOC 2 / industry-specific}}",
        "Incident response: {{plan in place}}",
    ], notes="Q: 'What about our data going to OpenAI?' A: '{{Specific PII redaction approach}}'")

    add("Quality measurement", "How we know it's working", [
        "Leading metric 1: {{e.g., time saved per transaction}}",
        "Leading metric 2: {{e.g., adoption %}}",
        "Lagging metric: {{e.g., customer NPS in {{period}}}}",
        "Quality bar: {{specific accuracy threshold}}",
        "Review cadence: {{weekly/monthly}}",
    ])

    add("Year 2 vision", "What this unlocks", [
        "Capacity for {{strategic project that was off the table}}",
        "Insights from {{data we'll have that we don't today}}",
        "Capability for {{new offering / market}}",
        "Foundation for {{next 10 automations}}",
    ], notes="Don't undersell. The first automation is just the foundation.")

    add("What we need from you", "The ask", [
        "Budget: ${{amount}} for {{period}}",
        "Headcount: {{people}} for {{period}}",
        "Decision: {{specific yes/no question}}",
        "Timeline: {{when}}",
        "Sponsor: {{role}} as executive sponsor",
    ], notes="Be explicit. Don't make them guess what you need.")

    add("Next steps", None, [
        "Week 1: {{action with owner}}",
        "Week 2: {{action with owner}}",
        "Week 4: {{milestone}}",
        "Week 12: {{first measurable outcome}}",
    ])

    add("Appendix · Full ROI breakdown", "All numbers", [
        "Inputs sheet (from ROI Calculator)",
        "12-month cash flow projection",
        "Sensitivity table — best / base / worst",
        "Comparison vs alternative approaches",
    ])

    add("Appendix · Why we haven't fixed this before", None, [
        "Previous attempt 1: {{what we tried, what went wrong}}",
        "Previous attempt 2: {{what we tried, what went wrong}}",
        "What's different now: {{AI capability that wasn't there 18 months ago}}",
        "What we learned: {{insights from previous attempts}}",
    ])

    add("Appendix · References", None, [
        "Case studies from Aiprosol (aiprosol.com/case-studies)",
        "Vendor security documentation (SOC 2, GDPR)",
        "Internal champion contacts willing to talk",
        "Industry benchmarks ({{source}})",
    ])

    add("Appendix · Glossary", None, [
        "HITL = Human-in-the-loop",
        "LLM = Large Language Model",
        "RAG = Retrieval-Augmented Generation",
        "TTV = Time to Value",
        "NPV = Net Present Value (discounted future cash flow)",
    ])

    add("Thank you", "Questions?", [
        "{{Your Name}}",
        "{{Your Email}}",
        "{{Your Phone}}",
        "{{Your Company}}",
    ])

    out = PD_DIR / "automation-roi-pitch-deck-template.pptx"
    prs.save(out)
    return out


pitch_pptx = build_pitch_deck()
print(f"✓ ROI Pitch Deck: {pitch_pptx.stat().st_size/1024:.1f} KB  {pitch_pptx.name}")


# ════════════════════════════════════════════════════════════════════
# 5. AUDIT CHECKLIST — Interactive HTML
# ════════════════════════════════════════════════════════════════════

AC_DIR = CATALOG / "checklists/037-business-process-audit-checklist/delivery"
AC_DIR.mkdir(parents=True, exist_ok=True)

CHECKLIST_DIMENSIONS = [
    ("People", [
        "Are there manual hand-offs between people for routine work?",
        "Do team members copy-paste data between systems?",
        "Are people doing the same task repeatedly (>5 times/week)?",
        "Do new hires need to learn 'tribal knowledge' that's not documented?",
        "Are senior staff doing tasks juniors could do?",
        "Are there single points of failure (one person who knows X)?",
        "Are people delayed waiting for approvals?",
        "Are people doing data entry from emails/PDFs?",
        "Are people manually building the same reports each week?",
        "Are people answering the same questions repeatedly?",
    ]),
    ("Process", [
        "Are there clear decision points (X happens, then Y)?",
        "Are exceptions well-documented?",
        "Are bottlenecks identified?",
        "Are cycle times measured?",
        "Are dependencies between processes mapped?",
        "Are SLAs defined for key steps?",
        "Are error/rework rates tracked?",
        "Is the happy path well understood?",
        "Are edge cases handled consistently?",
        "Are process owners identified?",
    ]),
    ("Data", [
        "Is data stored in multiple places that should be in one?",
        "Are there discrepancies between systems (same field, different values)?",
        "Is data freshness an issue?",
        "Are data quality issues frustrating people?",
        "Is there a single source of truth for key entities?",
        "Are reports built from manual data assembly?",
        "Is sensitive data well-controlled?",
        "Are integrations stable?",
        "Are exports/imports manual?",
        "Are there 'shadow' spreadsheets that should be in a system?",
    ]),
    ("Tools", [
        "Are people using 5+ tools daily?",
        "Are similar tools used redundantly?",
        "Are integrations between tools well-established?",
        "Are tools being underused (paying for features not used)?",
        "Is tool sprawl impacting productivity?",
        "Are tools chosen ad-hoc vs strategically?",
        "Are licenses reconciled against actual usage?",
        "Are there 'cohort' tools (everyone uses same)?",
        "Are tool changes coordinated with teams?",
        "Are tool gotchas documented?",
    ]),
    ("Outcomes", [
        "Are key metrics measured weekly?",
        "Are leading indicators visible?",
        "Are outcomes tied to compensation/recognition?",
        "Are dashboards built from current data (not 1-month-old)?",
        "Are anomalies surfaced automatically?",
        "Are trends (vs point-in-time) visible?",
        "Are forecasts vs actuals tracked?",
        "Are wins celebrated visibly?",
        "Are failures debriefed without blame?",
        "Is the team aligned on what 'success' looks like?",
    ]),
]


def build_audit_checklist():
    questions_html = []
    q_id = 1
    for dim_name, qs in CHECKLIST_DIMENSIONS:
        questions_html.append(f'<section class="dim" data-dim="{dim_name}"><h2>{dim_name}</h2>')
        for q in qs:
            questions_html.append(f'''
        <div class="q">
          <label class="qtext">{q}</label>
          <div class="opts">
            <label><input type="radio" name="q{q_id}" value="0" data-dim="{dim_name}">Never</label>
            <label><input type="radio" name="q{q_id}" value="1" data-dim="{dim_name}">Sometimes</label>
            <label><input type="radio" name="q{q_id}" value="2" data-dim="{dim_name}">Often</label>
            <label><input type="radio" name="q{q_id}" value="3" data-dim="{dim_name}">Always</label>
          </div>
        </div>''')
            q_id += 1
        questions_html.append('</section>')

    html = '''<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8"><meta name="viewport" content="width=device-width, initial-scale=1"><title>Business Process Audit Checklist — Aiprosol</title>
<style>
* { box-sizing: border-box; }
body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Inter, system-ui, sans-serif; max-width: 880px; margin: 0 auto; padding: 56px 32px 96px; line-height: 1.65; color: #1F1B2E; background: #FAF8FF; }
@media (prefers-color-scheme: dark) { body { color: #E5E7EB; background: #0A0613; } .q, .opts label { background: #1A1428; } input[type="text"], textarea { background: #1A1428; color: inherit; border-color: #2A1F3D; } .opts label { color: #C7CEDB; } }
h1 { font-size: 36px; margin: 0 0 8px; font-weight: 800; background: linear-gradient(135deg, #8B5CF6, #C084FC); -webkit-background-clip: text; -webkit-text-fill-color: transparent; background-clip: text; }
h2 { font-size: 22px; margin-top: 40px; margin-bottom: 12px; color: #8B5CF6; border-bottom: 1px solid rgba(139,92,246,0.3); padding-bottom: 6px; }
.q { background: rgba(139,92,246,0.04); padding: 14px 20px; margin: 14px 0; border-radius: 8px; }
.qtext { font-weight: 500; display: block; margin-bottom: 8px; }
.opts { display: flex; flex-wrap: wrap; gap: 12px; }
.opts label { background: white; padding: 6px 14px; border-radius: 999px; cursor: pointer; font-size: 14px; border: 1px solid rgba(139,92,246,0.2); }
.opts label:hover { border-color: #8B5CF6; }
.opts input { margin-right: 6px; }
.actions { margin: 40px 0; display: flex; gap: 12px; flex-wrap: wrap; }
button { background: linear-gradient(135deg, #8B5CF6, #C084FC); color: white; border: none; padding: 12px 24px; border-radius: 8px; font-size: 16px; font-weight: 600; cursor: pointer; font-family: inherit; }
button:hover { opacity: 0.9; }
button.secondary { background: transparent; color: #8B5CF6; border: 1px solid #8B5CF6; }
.results { margin: 40px 0; padding: 24px; background: rgba(139,92,246,0.06); border-radius: 12px; border-left: 4px solid #8B5CF6; }
.score-row { display: flex; justify-content: space-between; align-items: center; margin: 10px 0; padding: 8px 0; border-bottom: 1px solid rgba(139,92,246,0.15); }
.score-bar { flex: 1; height: 18px; background: rgba(139,92,246,0.1); border-radius: 8px; margin: 0 14px; overflow: hidden; position: relative; }
.score-bar-fill { height: 100%; background: linear-gradient(90deg, #EF4444, #F59E0B, #22C55E); border-radius: 8px; }
.tier-badge { display: inline-block; padding: 3px 10px; border-radius: 999px; font-size: 12px; font-weight: 700; text-transform: uppercase; letter-spacing: 0.05em; }
.tier-high { background: #FEE2E2; color: #991B1B; }
.tier-med { background: #FEF3C7; color: #92400E; }
.tier-low { background: #D1FAE5; color: #065F46; }
.opportunity-map { margin: 24px 0; }
.opp-row { background: white; padding: 14px 18px; margin: 8px 0; border-radius: 8px; display: flex; justify-content: space-between; gap: 18px; }
@media (prefers-color-scheme: dark) { .opp-row { background: #1A1428; } }
textarea { width: 100%; padding: 10px; border: 1px solid rgba(139,92,246,0.3); border-radius: 6px; font-family: inherit; min-height: 60px; margin-top: 8px; }
.brand-footer { margin-top: 60px; padding-top: 28px; border-top: 1px solid rgba(139,92,246,0.18); font-size: 12px; color: #6B7585; text-align: center; }
</style></head><body>

<h1>Business Process Audit Checklist</h1>
<p style="color: #6B7585; margin: 0 0 28px;">Find your highest-ROI automation candidates in 90 minutes. 50 questions across 5 dimensions. © 2026 Aiprosol Ltd</p>

<p>For each question, pick the answer that's most true today. Don't overthink. Don't gild it. Submit at the bottom and see your scored opportunity map.</p>

<form id="auditForm">
''' + "\n".join(questions_html) + '''

<section>
<h2>Notes (optional)</h2>
<textarea id="notes" placeholder="Any context you want to capture..."></textarea>
</section>

<div class="actions">
<button type="button" onclick="scoreAudit()">Score this audit</button>
<button type="button" class="secondary" onclick="resetAudit()">Reset</button>
<button type="button" class="secondary" onclick="exportResults()">Export results (JSON)</button>
</div>
</form>

<div id="results" class="results" style="display:none;">
<h2 style="margin-top:0;">Your audit results</h2>
<div id="dim-scores"></div>
<h3>Top automation opportunities (highest score = most upside)</h3>
<div id="opp-map" class="opportunity-map"></div>
<p style="margin-top: 20px; font-size: 14px; color: #6B7585;">For each top dimension, the corresponding Aiprosol product covers the methodology + n8n workflows you can deploy this week. See aiprosol.com/products.</p>
</div>

<div class="brand-footer">© 2026 Aiprosol Ltd · <a href="https://aiprosol.com" style="color:#8B5CF6;">aiprosol.com</a> · <a href="mailto:hello@aiprosol.com" style="color:#8B5CF6;">hello@aiprosol.com</a><br>Licensed to the purchaser for unlimited internal use.</div>

<script>
const DIM_PRODUCT_MAP = {
  "People": "Workflow Automation Playbook ($97) — 25 n8n workflows including AI-assisted handoffs",
  "Process": "30-Day Business Automation Challenge ($47) — 30 daily exercises for systematic process automation",
  "Data": "AI Tools Stack Starter Kit ($197) — 18 integration recipes connecting data across tools",
  "Tools": "AI Tools Vault ($147) — 96 tools curated with verdicts + hidden gems + avoid list",
  "Outcomes": "Lead Generation Automation Playbook ($127) — outcome-driven measurement + 12 n8n workflows"
};

function scoreAudit() {
  const dims = {};
  document.querySelectorAll('input[type="radio"]:checked').forEach(r => {
    const d = r.dataset.dim;
    dims[d] = dims[d] || { sum: 0, count: 0 };
    dims[d].sum += Number(r.value);
    dims[d].count++;
  });

  // Score per dim: average × 100/3
  const scores = Object.entries(dims).map(([dim, s]) => ({
    dim,
    score: Math.round(s.sum / s.count * 100 / 3) || 0,
    answered: s.count
  })).sort((a, b) => b.score - a.score);

  const container = document.getElementById('dim-scores');
  container.innerHTML = scores.map(s => {
    let tier = 'tier-low';
    let tierLabel = 'Low priority';
    if (s.score >= 65) { tier = 'tier-high'; tierLabel = 'High priority'; }
    else if (s.score >= 35) { tier = 'tier-med'; tierLabel = 'Medium priority'; }
    return `<div class="score-row">
      <strong style="min-width:100px;">${s.dim}</strong>
      <div class="score-bar"><div class="score-bar-fill" style="width:${s.score}%"></div></div>
      <span style="min-width:50px;text-align:right;font-weight:600;">${s.score}/100</span>
      <span class="tier-badge ${tier}">${tierLabel}</span>
    </div>`;
  }).join('');

  const opps = scores.slice(0, 3);
  const oppContainer = document.getElementById('opp-map');
  oppContainer.innerHTML = opps.map((s, i) => `<div class="opp-row">
    <div>
      <strong>${i+1}. ${s.dim}</strong> (score ${s.score}/100)<br>
      <span style="font-size:13px;color:#6B7585;">${DIM_PRODUCT_MAP[s.dim] || ''}</span>
    </div>
  </div>`).join('');

  document.getElementById('results').style.display = 'block';
  document.getElementById('results').scrollIntoView({behavior:'smooth'});
}

function resetAudit() {
  document.querySelectorAll('input[type="radio"]').forEach(r => r.checked = false);
  document.getElementById('notes').value = '';
  document.getElementById('results').style.display = 'none';
}

function exportResults() {
  const dims = {};
  document.querySelectorAll('input[type="radio"]:checked').forEach(r => {
    const d = r.dataset.dim;
    dims[d] = dims[d] || { sum: 0, count: 0 };
    dims[d].sum += Number(r.value);
    dims[d].count++;
  });
  const scores = Object.entries(dims).map(([dim, s]) => ({dim, score: Math.round(s.sum/s.count*100/3)||0}));
  const data = {scores, notes: document.getElementById('notes').value, date: new Date().toISOString()};
  const blob = new Blob([JSON.stringify(data, null, 2)], {type:'application/json'});
  const url = URL.createObjectURL(blob);
  const a = document.createElement('a');
  a.href = url; a.download = `audit-${new Date().toISOString().slice(0,10)}.json`; a.click();
}
</script>
</body></html>'''

    out = AC_DIR / "business-process-audit-checklist.html"
    out.write_text(html)
    return out


checklist_html = build_audit_checklist()
print(f"✓ Audit Checklist: {checklist_html.stat().st_size/1024:.1f} KB  {checklist_html.name}")

print("\nAll 5 products generated.")
