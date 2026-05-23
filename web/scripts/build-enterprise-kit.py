#!/usr/bin/env python3
"""
Build the 3 binary deliverables for Enterprise AI Readiness Assessment Kit:
  1. readiness-scorecard.xlsx  — 120 questions × 12 dimensions, auto-scoring
  2. vendor-rfp-template.docx  — 40-criteria evaluation grid + red-flag screen
  3. 90-day-plan.pptx          — 14-slide implementation deck

Brand: violet (#8B5CF6 + #C084FC), Space Grotesk for headlines.
"""

from pathlib import Path
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, Reference
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).resolve().parent.parent / "private-products"

# ─────────────────────────────────────────────────────────────────────────
# 1. READINESS SCORECARD — 120 questions × 12 dimensions
# ─────────────────────────────────────────────────────────────────────────

DIMENSIONS = [
    ("Data foundation", [
        "Do we have a unified customer record across all systems?",
        "Are the top-3 datasets queryable by non-engineers?",
        "Is data lineage documented for revenue, customer, product data?",
        "Do we have a primary source-of-truth for customer identity?",
        "Are data refresh latencies known and acceptable (<24h for ops)?",
        "Do we have data quality checks running on critical pipelines?",
        "Is there a documented retention policy for customer data?",
        "Are PII fields tagged and access-controlled?",
        "Do we have a data dictionary covering top-20 tables?",
        "Are data exports auditable (who pulled what, when)?",
    ]),
    ("Identity + access", [
        "Do we have SSO across all production systems?",
        "Are role-based access controls audited quarterly?",
        "Do we have a documented offboarding checklist that completes in <24h?",
        "Is MFA enforced for all admin-level access?",
        "Do we use a password manager organisation-wide?",
        "Are service accounts inventoried with named owners?",
        "Do we rotate API keys on a documented schedule?",
        "Is there a process for emergency access revocation?",
        "Are guest/contractor accounts time-bound by default?",
        "Do we audit account access logs at least monthly?",
    ]),
    ("Tool architecture", [
        "Do we know exactly how many SaaS tools we operate?",
        "Are integrations hub-and-spoke vs point-to-point?",
        "Is there a workflow orchestrator already in place (n8n / Zapier / Make)?",
        "Have we identified our top-5 tool overlaps (redundancy)?",
        "Do we have a SaaS spend budget owner?",
        "Are tool renewals tracked centrally?",
        "Is there a process for proposing new tools?",
        "Do we sunset unused tools quarterly?",
        "Are vendor SOC2 reports collected and reviewed?",
        "Do we have a documented integration architecture diagram?",
    ]),
    ("AI literacy", [
        "What % of staff have used a frontier LLM in the last 30 days? (≥50% = score full)",
        "Do we have an internal Slack channel for AI use cases?",
        "Is there a designated AI owner / champion?",
        "Have we run an AI training session in the last 6 months?",
        "Do we maintain an internal prompt library?",
        "Can leadership demonstrate hands-on AI tool use?",
        "Do we have an AI use-case backlog?",
        "Is there a recurring AI show-and-tell cadence?",
        "Do staff have access to paid AI tool licenses where needed?",
        "Have we measured AI usage impact on at least one workflow?",
    ]),
    ("Governance + risk", [
        "Do we have a written AI usage policy?",
        "Is there a process for approving new AI tools?",
        "Are AI-generated outputs labelled as such?",
        "Do we have a customer-disclosure stance on AI use?",
        "Is there a board-level review of AI initiatives?",
        "Do we have insurance covering AI-related liability?",
        "Are AI vendor contracts reviewed by legal before signing?",
        "Do we maintain an AI incident response playbook?",
        "Is bias / fairness testing applied to customer-facing AI?",
        "Do we have an opt-out mechanism for customer data in AI training?",
    ]),
    ("Compliance posture", [
        "Are we GDPR compliant (if EU customers)?",
        "Are we CCPA compliant (if California customers)?",
        "Do we have SOC 2 Type 2 (or active path)?",
        "Are AI processors covered by data processing agreements (DPAs)?",
        "Do we maintain data residency where required?",
        "Are AI-generated decisions documented for audit?",
        "Do we have a documented vendor security review process?",
        "Are penetration tests run at least annually?",
        "Do we maintain a security incident response plan?",
        "Are customer DPAs updated to cover AI sub-processors?",
    ]),
    ("Talent", [
        "Do we have someone who can specify an AI workflow?",
        "Do we have someone who can build it (code or no-code)?",
        "Do we have someone who can operate it day-to-day?",
        "Do we have someone who can audit it for drift?",
        "Have we budgeted for AI-related training in the next 12 months?",
        "Are job descriptions updated to include AI fluency expectations?",
        "Do we have access to AI consulting for gaps?",
        "Is there a career path for AI specialists internally?",
        "Have we mapped which roles are most automatable?",
        "Do we have a plan for displaced or augmented roles?",
    ]),
    ("Process maturity", [
        "Are our top-5 business processes documented?",
        "Are process SLAs measured?",
        "Is there a recurring process improvement cadence (monthly+)?",
        "Do we have process owners named for each top-5 process?",
        "Are exceptions (edge cases) documented per process?",
        "Do we maintain runbooks for critical workflows?",
        "Are processes mapped to underlying tools used?",
        "Have we measured task-frequency × duration for top processes?",
        "Are process changes communicated to all stakeholders?",
        "Do we sunset broken processes systematically?",
    ]),
    ("Change management", [
        "How does the org typically respond to new tools? (positive = full score)",
        "Who owns adoption when a new tool launches?",
        "Is there a feedback loop after rollouts?",
        "Do we measure adoption rate (not just deployment)?",
        "Are change announcements consistent in format?",
        "Do we provide training before rollouts?",
        "Are change champions identified per team?",
        "Do we run pilots before full rollout?",
        "Are rollback plans documented?",
        "Do we celebrate successful adoption wins?",
    ]),
    ("Financial model", [
        "Can we measure ROI of automation projects?",
        "Is there a budget line for AI tooling?",
        "Are we tracking task-cost vs fully-loaded-employee-cost?",
        "Do we have a payback-period threshold for automation projects?",
        "Are tool costs allocated to business units accurately?",
        "Do we forecast AI cost growth annually?",
        "Is there a cost-per-workflow metric tracked?",
        "Are AI token costs monitored and budgeted?",
        "Do we benchmark AI spend against industry peers?",
        "Is there a CFO-level review of AI spend?",
    ]),
    ("Customer impact", [
        "Are customer-facing AI features labelled?",
        "Is there a customer feedback channel for AI outputs?",
        "Do we A/B test AI vs human alternatives?",
        "Have we surveyed customers on AI comfort levels?",
        "Are escalation paths to humans clear?",
        "Do we measure customer satisfaction post-AI interaction?",
        "Are AI errors tracked from a customer-impact angle?",
        "Do we have a public AI principles statement?",
        "Are customers given control over AI use on their data?",
        "Do we measure trust scores related to AI features?",
    ]),
    ("Strategic clarity", [
        "Do we have a 12-month AI roadmap?",
        "Is the board briefed on AI quarterly?",
        "Is there a defensible thesis for why AI is core to our business?",
        "Have we identified AI moats vs commodity AI features?",
        "Is AI a top-3 strategic priority for the year?",
        "Do we know our 'AI competitive position' vs peers?",
        "Is there a 3-year AI investment thesis?",
        "Are AI initiatives tied to strategic OKRs?",
        "Do we have an AI failure tolerance / risk appetite?",
        "Have we communicated AI strategy to all-hands?",
    ]),
]


def build_scorecard():
    wb = Workbook()
    ws = wb.active
    ws.title = "Readiness Scorecard"

    # Brand colors
    violet = "8B5CF6"
    violet_light = "C084FC"
    dark = "1F1B2E"
    bg_alt = "F4F1FA"
    border_light = Border(
        left=Side(style="thin", color="D4D4D4"),
        right=Side(style="thin", color="D4D4D4"),
        top=Side(style="thin", color="D4D4D4"),
        bottom=Side(style="thin", color="D4D4D4"),
    )

    # Title row
    ws["A1"] = "Aiprosol · Enterprise AI Readiness Scorecard"
    ws["A1"].font = Font(name="Calibri", size=18, bold=True, color=violet)
    ws.merge_cells("A1:E1")
    ws.row_dimensions[1].height = 32

    ws["A2"] = "Score each question 0–10 (0 = not at all, 5 = partial, 10 = fully). Each dimension max = 100."
    ws["A2"].font = Font(name="Calibri", size=10, italic=True, color="666666")
    ws.merge_cells("A2:E2")

    # Headers
    headers = ["Dimension", "#", "Question", "Score 0–10", "Notes"]
    row = 4
    for col, h in enumerate(headers, start=1):
        c = ws.cell(row=row, column=col, value=h)
        c.font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor=violet)
        c.alignment = Alignment(horizontal="left", vertical="center")
        c.border = border_light

    ws.column_dimensions["A"].width = 22
    ws.column_dimensions["B"].width = 5
    ws.column_dimensions["C"].width = 70
    ws.column_dimensions["D"].width = 12
    ws.column_dimensions["E"].width = 40

    row += 1
    dim_start_rows = {}
    for dim_name, questions in DIMENSIONS:
        dim_start_rows[dim_name] = row
        for i, q in enumerate(questions, start=1):
            ws.cell(row=row, column=1, value=dim_name if i == 1 else "").font = Font(bold=(i == 1), color=violet if i == 1 else "000000")
            ws.cell(row=row, column=2, value=i)
            ws.cell(row=row, column=3, value=q)
            ws.cell(row=row, column=4, value=None)  # user fills
            ws.cell(row=row, column=5, value="")
            for col in range(1, 6):
                ws.cell(row=row, column=col).alignment = Alignment(vertical="center", wrap_text=(col == 3))
                ws.cell(row=row, column=col).border = border_light
                if (row - 4) % 2 == 0:
                    ws.cell(row=row, column=col).fill = PatternFill("solid", fgColor=bg_alt)
            ws.row_dimensions[row].height = 22
            row += 1

    # Summary section
    summary_start = row + 2
    ws.cell(row=summary_start, column=1, value="DIMENSION SUMMARY").font = Font(size=14, bold=True, color=violet)
    ws.merge_cells(start_row=summary_start, start_column=1, end_row=summary_start, end_column=5)
    summary_start += 1

    summary_headers = ["Dimension", "Sub-total (max 100)", "% Complete", "Tier", "Priority"]
    for col, h in enumerate(summary_headers, start=1):
        c = ws.cell(row=summary_start, column=col, value=h)
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor=violet)
        c.border = border_light

    summary_start += 1
    dim_summary_first = summary_start
    for dim_name, _qs in DIMENSIONS:
        start = dim_start_rows[dim_name]
        end = start + 9
        ws.cell(row=summary_start, column=1, value=dim_name).font = Font(bold=True)
        ws.cell(row=summary_start, column=2, value=f"=SUM(D{start}:D{end})")
        ws.cell(row=summary_start, column=3, value=f"=B{summary_start}/100").number_format = "0%"
        ws.cell(row=summary_start, column=4, value=f'=IF(C{summary_start}>=0.85,"Leading",IF(C{summary_start}>=0.7,"Advanced",IF(C{summary_start}>=0.5,"Operating",IF(C{summary_start}>=0.3,"Building","Foundational"))))')
        ws.cell(row=summary_start, column=5, value=f'=IF(C{summary_start}<0.5,"HIGH","")')
        for col in range(1, 6):
            ws.cell(row=summary_start, column=col).border = border_light
        summary_start += 1
    dim_summary_last = summary_start - 1

    # Overall total
    summary_start += 1
    ws.cell(row=summary_start, column=1, value="OVERALL").font = Font(size=14, bold=True, color=violet)
    ws.cell(row=summary_start, column=2, value=f"=SUM(B{dim_summary_first}:B{dim_summary_last})").font = Font(size=14, bold=True)
    ws.cell(row=summary_start, column=3, value=f"=B{summary_start}/1200").number_format = "0%"
    ws.cell(row=summary_start, column=3).font = Font(size=14, bold=True)
    ws.cell(row=summary_start, column=4, value=f'=IF(C{summary_start}>=0.85,"LEADING",IF(C{summary_start}>=0.7,"ADVANCED",IF(C{summary_start}>=0.5,"OPERATING",IF(C{summary_start}>=0.3,"BUILDING","FOUNDATIONAL"))))').font = Font(size=14, bold=True, color=violet)

    # Tier interpretation legend
    summary_start += 3
    ws.cell(row=summary_start, column=1, value="TIER GUIDE").font = Font(size=12, bold=True, color=violet)
    summary_start += 1
    tiers = [
        ("0–30%  Foundational", "Hire/contract an AI lead first. Don't roll out enterprise AI yet."),
        ("31–50% Building", "Pick 1–2 high-leverage workflows. Pilot for 90 days."),
        ("51–70% Operating", "Scale pilots to 5–10 production workflows. Establish governance."),
        ("71–85% Advanced", "Compound: center of excellence. Train internal champions."),
        ("86–100% Leading", "Ahead of 95% of peers. Focus on talent retention + safety."),
    ]
    for tier, advice in tiers:
        ws.cell(row=summary_start, column=1, value=tier).font = Font(bold=True)
        ws.cell(row=summary_start, column=3, value=advice).alignment = Alignment(wrap_text=True)
        ws.row_dimensions[summary_start].height = 24
        summary_start += 1

    # Footer
    summary_start += 2
    ws.cell(row=summary_start, column=1, value="Aiprosol · aiprosol.com · Redeem your 30-min scoring call: srijanpaudelofficial@gmail.com").font = Font(italic=True, size=9, color="888888")

    # Chart for dimension scores
    try:
        chart = BarChart()
        chart.type = "bar"
        chart.style = 12
        chart.title = "Readiness by Dimension"
        chart.y_axis.title = "Dimension"
        chart.x_axis.title = "Score (0–100)"
        data = Reference(ws, min_col=2, min_row=dim_summary_first - 1, max_row=dim_summary_last, max_col=2)
        cats = Reference(ws, min_col=1, min_row=dim_summary_first, max_row=dim_summary_last)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.height = 12
        chart.width = 18
        ws.add_chart(chart, "G4")
    except Exception as e:
        print(f"  chart skip: {e}")

    out = OUT / "enterprise-readiness-scorecard.xlsx"
    wb.save(out)
    return out


# ─────────────────────────────────────────────────────────────────────────
# 2. VENDOR RFP TEMPLATE — DOCX
# ─────────────────────────────────────────────────────────────────────────

RFP_CATEGORIES = [
    ("Technical fit", [
        "Integration footprint with our existing stack (HubSpot, Stripe, Slack, Notion, Google Workspace)",
        "REST API quality: documented OpenAPI spec, predictable response times, semantic versioning",
        "Latency SLAs: median + p95 + p99 published",
        "Model choice flexibility: support for Claude, GPT, Llama (avoiding single-vendor lock-in)",
        "Fine-tuning / custom model support",
        "On-premises or VPC deployment option",
        "Audit logs accessible via API, retained 90+ days",
        "Data residency options for our customer base",
    ]),
    ("Commercial terms", [
        "Pricing model: per-seat / per-API-call / flat / hybrid",
        "Contract lock-in: minimum term and renewal terms",
        "Termination terms: notice period, data retrieval guarantees",
        "SLA credits for downtime",
        "Professional services / onboarding hours included",
    ]),
    ("Security & compliance", [
        "SOC 2 Type 2 report (within last 12 months)",
        "ISO 27001 certification",
        "GDPR-compliant DPA template available",
        "Encryption at rest (AES-256) and in transit (TLS 1.2+)",
        "Vulnerability disclosure / responsible disclosure policy",
    ]),
    ("Customer support", [
        "Tiered support availability (24/7, business hours, basic)",
        "Response SLAs by severity",
        "Named account manager (paid tiers)",
        "Dedicated Slack channel or shared support inbox",
        "Public status page with incident history",
    ]),
    ("Roadmap alignment", [
        "Public roadmap or quarterly roadmap reviews",
        "Feature request process with transparency",
        "Customer advisory board access",
        "Release cadence (minor / major versions)",
    ]),
    ("Reference customers", [
        "References in our industry willing to speak",
        "Reference customers of similar size (10–500 employees)",
        "Case studies with measurable outcomes (ROI, hours reclaimed, etc.)",
        "Publicly listed customers we can validate independently",
    ]),
    ("Implementation support", [
        "Onboarding sprint with named milestones",
        "Dedicated Customer Success Manager",
        "Training resources (docs, videos, live sessions) included",
        "Sandbox environment available for evaluation",
        "Migration assistance from incumbent systems",
    ]),
    ("Long-term viability", [
        "Funding runway / profitability status",
        "Parent company / acquisition history",
        "Customer count trajectory (last 24 months)",
        "Employee Glassdoor signal + team retention",
    ]),
]

RED_FLAGS = [
    "Can you provide a SOC 2 Type 2 report dated within the last 12 months?",
    "Will you sign our standard DPA without material modifications?",
    "Will you commit in writing to delete our data within 30 days of termination?",
    "Can you provide three reference customers in our industry, willing to speak by phone?",
    "What is your customer churn rate over the last 12 months?",
    "Have you been acquired in the last 12 months OR are you actively raising capital this quarter?",
    "Will you indemnify us against IP infringement claims arising from your training data?",
    "Do you guarantee 30+ days written notice before breaking API changes?",
    "Can we run a 30-day technical evaluation in a sandbox before signing?",
    "Will you commit to a maximum 30-day onboarding window with named milestones and SLAs?",
]


def build_rfp():
    doc = Document()

    # Title
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title.add_run("AI Vendor RFP & Evaluation Template")
    title_run.font.name = "Calibri"
    title_run.font.size = Pt(22)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x8B, 0x5C, 0xF6)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = sub.add_run("Aiprosol Enterprise Kit · v1.0 · 2026")
    sub_run.font.size = Pt(10)
    sub_run.font.italic = True
    sub_run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    doc.add_paragraph()
    doc.add_paragraph(
        "Use this template to evaluate AI vendors consistently. Each criterion is scored 1–5 by your evaluation committee. "
        "Total max score: 200. Vendors below 140 should be deprioritised. Vendors below 100 should be eliminated."
    )

    doc.add_heading("Vendor Information", level=1)
    info_table = doc.add_table(rows=6, cols=2)
    info_table.style = "Light Grid Accent 4"
    fields = [
        "Vendor name:",
        "Primary contact name + email:",
        "Date of RFP response:",
        "Evaluation committee members:",
        "Target go-live date:",
        "Budget envelope (annual):",
    ]
    for i, label in enumerate(fields):
        info_table.cell(i, 0).text = label
        info_table.cell(i, 0).paragraphs[0].runs[0].font.bold = True

    # Evaluation grid
    doc.add_heading("40-Criteria Evaluation Grid", level=1)
    doc.add_paragraph("Score each row 1–5: (1=poor, 2=below average, 3=meets, 4=exceeds, 5=best in class)")

    for cat_name, items in RFP_CATEGORIES:
        doc.add_heading(cat_name, level=2)
        table = doc.add_table(rows=len(items) + 1, cols=4)
        table.style = "Light Grid Accent 4"
        hdr_cells = table.rows[0].cells
        for i, h in enumerate(["#", "Criterion", "Score (1–5)", "Notes"]):
            hdr_cells[i].text = h
            hdr_cells[i].paragraphs[0].runs[0].font.bold = True
        for i, item in enumerate(items, start=1):
            row = table.rows[i].cells
            row[0].text = str(i)
            row[1].text = item
            row[2].text = ""
            row[3].text = ""

    # Red flags
    doc.add_heading("Red-Flag Screening Questions (Yes/No)", level=1)
    doc.add_paragraph(
        "Each vendor must answer these 10 in writing. Any 'No' or evasive answer is grounds to eliminate from shortlist."
    )
    rf_table = doc.add_table(rows=len(RED_FLAGS) + 1, cols=3)
    rf_table.style = "Light Grid Accent 4"
    rf_hdr = rf_table.rows[0].cells
    for i, h in enumerate(["#", "Question", "Answer (Y/N + comment)"]):
        rf_hdr[i].text = h
        rf_hdr[i].paragraphs[0].runs[0].font.bold = True
    for i, q in enumerate(RED_FLAGS, start=1):
        row = rf_table.rows[i].cells
        row[0].text = str(i)
        row[1].text = q
        row[2].text = ""

    # Decision summary
    doc.add_heading("Decision Summary", level=1)
    sum_table = doc.add_table(rows=5, cols=2)
    sum_table.style = "Light Grid Accent 4"
    sum_items = [
        ("Total evaluation score (out of 200):", ""),
        ("Red flags passed (out of 10):", ""),
        ("Recommended decision:", "Proceed / Shortlist / Decline"),
        ("Decision date:", ""),
        ("Signatories:", ""),
    ]
    for i, (label, val) in enumerate(sum_items):
        sum_table.cell(i, 0).text = label
        sum_table.cell(i, 0).paragraphs[0].runs[0].font.bold = True
        sum_table.cell(i, 1).text = val

    # Footer
    doc.add_paragraph()
    foot = doc.add_paragraph()
    foot.alignment = WD_ALIGN_PARAGRAPH.CENTER
    foot_run = foot.add_run("Aiprosol · aiprosol.com · 30-min scoring call: srijanpaudelofficial@gmail.com")
    foot_run.font.size = Pt(9)
    foot_run.font.italic = True
    foot_run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    out = OUT / "enterprise-vendor-rfp-template.docx"
    doc.save(out)
    return out


# ─────────────────────────────────────────────────────────────────────────
# 3. 90-DAY IMPLEMENTATION PLAN — PPTX (14 slides)
# ─────────────────────────────────────────────────────────────────────────

VIOLET = PptxRGBColor(0x8B, 0x5C, 0xF6)
VIOLET_LIGHT = PptxRGBColor(0xC0, 0x84, 0xFC)
DARK = PptxRGBColor(0x1F, 0x1B, 0x2E)
WHITE = PptxRGBColor(0xFF, 0xFF, 0xFF)
GREY = PptxRGBColor(0x80, 0x80, 0x80)


def add_slide(prs, title, subtitle=None, content_lines=None, footer="Aiprosol · Enterprise AI Readiness Kit"):
    blank = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank)

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(13.333), PptxInches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = VIOLET
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(0.5), PptxInches(12), PptxInches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(36)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = "Calibri"

    if subtitle:
        sub_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(1.5), PptxInches(12), PptxInches(0.6))
        sf = sub_box.text_frame
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = PptxPt(18)
        sp.font.color.rgb = VIOLET
        sp.font.italic = True

    if content_lines:
        body_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(2.3), PptxInches(12), PptxInches(4.5))
        bf = body_box.text_frame
        bf.word_wrap = True
        for i, line in enumerate(content_lines):
            p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            p.text = "•  " + line
            p.font.size = PptxPt(18)
            p.font.color.rgb = DARK
            p.space_after = PptxPt(10)

    # Footer
    foot_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(7.0), PptxInches(12), PptxInches(0.3))
    ff = foot_box.text_frame
    fp = ff.paragraphs[0]
    fp.text = footer
    fp.font.size = PptxPt(9)
    fp.font.color.rgb = GREY
    fp.font.italic = True

    return slide


def build_plan_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Slide 1 — Title
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(13.333), PptxInches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    title_box = s1.shapes.add_textbox(PptxInches(0.8), PptxInches(2.5), PptxInches(12), PptxInches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "90-Day AI Implementation Plan"
    p.font.size = PptxPt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = "Phase 0 → Phase 1 → Phase 2"
    p2.font.size = PptxPt(24)
    p2.font.color.rgb = VIOLET_LIGHT
    p2.space_before = PptxPt(20)
    foot_box = s1.shapes.add_textbox(PptxInches(0.8), PptxInches(6.7), PptxInches(12), PptxInches(0.5))
    ff = foot_box.text_frame
    fp = ff.paragraphs[0]
    fp.text = "Aiprosol · Enterprise AI Readiness Kit · 2026"
    fp.font.size = PptxPt(11)
    fp.font.color.rgb = VIOLET_LIGHT

    # Slide 2 — Summary KPIs
    add_slide(prs, "Summary KPIs", "Target outcomes over 90 days", [
        "Top-3 workflows live in production by day 60",
        "Average response time on key process reduced by 70%+",
        "Hours reclaimed per week: 30+ across pilot cohort",
        "100% of relevant staff trained by day 90",
        "Annual cost savings: $250k–$1M (mid-size org baseline)",
    ])

    # Slide 3 — Today's baseline
    add_slide(prs, "Today's Baseline", "From the Readiness Scorecard", [
        "Overall readiness score: [INSERT FROM SCORECARD] / 100",
        "Strongest dimensions: [TOP 3]",
        "Weakest dimensions: [BOTTOM 3 — these are our 90-day focus]",
        "Current SaaS spend on automation: $[INSERT] / mo",
        "Estimated automation candidates from audit: [INSERT] workflows",
    ])

    # Slide 4 — Selected vendors
    add_slide(prs, "Selected Vendors", "Output of RFP scoring", [
        "[Vendor 1] — primary platform — score [X]/200",
        "[Vendor 2] — secondary / specialised — score [X]/200",
        "Rationale: [why this combination]",
        "Contract terms: [duration, $$ committed, exit clauses]",
        "Implementation partner: Aiprosol (or other) for first 90 days",
    ])

    # Slide 5 — Phase 0 (Days 0–14)
    add_slide(prs, "Phase 0 — Foundations (Days 0–14)", "Set the stage for execution", [
        "Week 1: SSO wiring, DPA execution, security review",
        "Week 1: Data audit on top-3 datasets needed for pilot",
        "Week 2: Pilot user cohort selected (8–15 users across 2–3 teams)",
        "Week 2: Success criteria documented + baselined per workflow",
        "Output: production access · 5-min usage walkthrough complete · baselines logged",
    ])

    # Slide 6 — Phase 0 risks
    add_slide(prs, "Phase 0 Risks", "What can go wrong + mitigation", [
        "Risk: SSO integration delays · Mitigation: backup local accounts for first 2 weeks",
        "Risk: DPA legal back-and-forth · Mitigation: standard contract pre-approved",
        "Risk: pilot cohort drops out · Mitigation: 2x oversubscribe, recruit champions",
        "Risk: baseline data not available · Mitigation: instrument week 0",
    ])

    # Slide 7 — Phase 1 (Days 15–45)
    add_slide(prs, "Phase 1 — Pilot (Days 15–45)", "First 3 workflows in production", [
        "Week 3–4: Workflow 1 built + tested + shipped to pilot cohort",
        "Week 4–5: Workflow 2 built + shipped",
        "Week 5–6: Workflow 3 built + shipped",
        "Daily: Slack-pinged metrics to leadership channel",
        "Weekly: 30-min cohort interview cadence",
        "Output: 3 workflows running · hours-reclaimed metric trending · 5+ user testimonials",
    ])

    # Slide 8 — Phase 1 risks
    add_slide(prs, "Phase 1 Risks", "What can go wrong + mitigation", [
        "Risk: workflow latency higher than expected · Mitigation: caching layer + async patterns",
        "Risk: AI accuracy below threshold · Mitigation: human-in-the-loop fallback",
        "Risk: pilot users disengage · Mitigation: weekly office hours + champion model",
        "Risk: budget overrun · Mitigation: weekly cost reports + hard stop at 110%",
    ])

    # Slide 9 — Phase 2 (Days 46–90)
    add_slide(prs, "Phase 2 — Scale (Days 46–90)", "Roll out to wider org", [
        "Week 7–8: Onboard remaining users (training, runbook handoff)",
        "Week 9–10: Add workflows 4–10 based on pilot feedback",
        "Week 11–12: Governance launch — AI policy, training programme, escalation paths",
        "Week 13: 90-day review + decision on Phase 3 (year 1 roadmap)",
        "Output: 100% relevant staff onboarded · 10+ workflows live · monthly governance review cadence",
    ])

    # Slide 10 — Phase 2 risks
    add_slide(prs, "Phase 2 Risks", "What can go wrong + mitigation", [
        "Risk: change fatigue across org · Mitigation: phased rollout, no more than 2 changes/team/month",
        "Risk: governance theatre (rules without enforcement) · Mitigation: monthly audit cadence",
        "Risk: customer-facing AI errors · Mitigation: clear escalation + correction loop",
        "Risk: vendor scope creep · Mitigation: change control board",
    ])

    # Slide 11 — Talent plan
    add_slide(prs, "Talent Plan", "Roles needed across the 90 days", [
        "AI Workflow Specialist: 0.5 FTE (internal or Aiprosol-supplied)",
        "Implementation Engineer: 1.0 FTE during Phase 0–1",
        "Workflow Operator(s): 0.2 FTE per workflow ongoing",
        "Auditor / governance owner: 0.1 FTE ongoing",
        "Executive sponsor: 1 hour / week, day-90 review owner",
    ])

    # Slide 12 — Governance + audit
    add_slide(prs, "Governance + Audit Framework", "Established by Day 60", [
        "AI Usage Policy: published, signed by every staff member with AI access",
        "Vendor approval workflow: new tool proposals reviewed bi-weekly",
        "Output review: random 1% sampled monthly, accuracy tracked",
        "Incident response: documented runbook + 24h SLA on customer-facing issues",
        "Quarterly board review of AI initiatives + outcomes",
    ])

    # Slide 13 — Financial model
    add_slide(prs, "Financial Model", "Cost vs. reclaim across the 90 days", [
        "Phase 0 cost: licenses + setup ≈ $15–30k",
        "Phase 1 cost: implementation + integration ≈ $30–80k",
        "Phase 2 cost: training + governance launch ≈ $20–40k",
        "Total: $65–150k investment over 90 days",
        "Expected reclaim by Day 90: 30+ hours / week across pilot cohort",
        "12-month projected savings: $250k–$1M (depends on org size + scope)",
    ])

    # Slide 14 — Day-90 review template
    add_slide(prs, "Day-90 Review Template", "Decision points + next steps", [
        "Workflows live: [N] of 10 target",
        "Hours reclaimed / week: [X] vs target [Y]",
        "User adoption: [%] of intended cohort",
        "Cost vs budget: [%] of envelope used",
        "Top 3 wins · Top 3 lessons · Top 3 risks for next quarter",
        "Phase 3 decision: scale (yes / yes-with-changes / no — and what next)",
    ])

    out = OUT / "enterprise-90-day-implementation-plan.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    print("Building Enterprise Kit binaries...")
    f1 = build_scorecard()
    print(f"  ✓ {f1.name} ({f1.stat().st_size // 1024} KB)")
    f2 = build_rfp()
    print(f"  ✓ {f2.name} ({f2.stat().st_size // 1024} KB)")
    f3 = build_plan_pptx()
    print(f"  ✓ {f3.name} ({f3.stat().st_size // 1024} KB)")
    print("Done.")
